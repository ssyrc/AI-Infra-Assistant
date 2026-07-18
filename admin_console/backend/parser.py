"""
업로드된 매뉴얼 파일을 검색용 청크 리스트로 변환한다.

파일 유형별로 자연스러운 경계로 청킹한다:
- docx / pdf : Docling 마크다운 -> '#' 헤더(섹션) 단위 -> 길면 문단 재분할
- pptx       : 슬라이드 단위(슬라이드 = 하나의 청크 후보) -> 길면 문단 재분할
- txt / md   : 빈 줄(문단) 경계 -> 길면 재분할
- xlsx       : (별도 흐름) routers/manuals.py의 엑셀 컬럼 선택 커밋에서 처리

모든 청크 텍스트에는 cleaning.clean_text가 적용되어 HTML/제어문자/잡음이 제거된다.
폐쇄망에서는 Docling 모델 캐시를 사전 반입해야 한다(README 참고).
"""
import os
import re
from dataclasses import dataclass

from cleaning import clean_text, CleanOptions

MAX_CHUNK_CHARS = 1500

_converter = None


def _get_converter():
    """Docling은 무겁고 모델 다운로드가 필요하므로 실제로 문서를 파싱할 때 최초 1회만 로딩한다.
    dev 환경에서 DISABLE_DOCLING=1이면 로딩하지 않는다(txt/xlsx는 Docling 없이도 동작)."""
    global _converter
    if os.environ.get("DISABLE_DOCLING") == "1":
        raise RuntimeError(
            "이 개발 환경에서는 Docling이 비활성화되어 있습니다(docx/pptx/pdf). "
            "txt와 엑셀 컬럼선택 흐름은 dev에서도 동작합니다."
        )
    if _converter is None:
        from docling.document_converter import DocumentConverter
        _converter = DocumentConverter()
    return _converter


@dataclass
class ParsedChunk:
    section_title: str | None
    page_no: int | None
    chunk_text: str


SUPPORTED_EXTS = {".docx", ".pptx", ".pdf", ".txt", ".md"}


def parse_file(filepath: str, opts: CleanOptions | None = None,
               include_speaker_notes: bool = False) -> list[ParsedChunk]:
    """확장자에 따라 알맞은 파서로 분기한다. opts로 정제 강도를 지정한다."""
    opts = opts or CleanOptions()
    ext = os.path.splitext(filepath)[1].lower()

    if ext in (".txt", ".md"):
        raw_sections = _parse_text_file(filepath)
    elif ext == ".pptx":
        raw_sections = _parse_pptx(filepath, include_speaker_notes)
    elif ext in (".docx", ".pdf"):
        raw_sections = _parse_via_docling(filepath)
    else:
        raise ValueError(f"지원하지 않는 확장자입니다: {ext}")

    chunks: list[ParsedChunk] = []
    for title, page_no, body in raw_sections:
        cleaned = clean_text(body, opts)
        if not cleaned:
            continue
        clean_title = clean_text(title, opts) if title else None
        for piece in _split_long_text(cleaned, MAX_CHUNK_CHARS):
            chunks.append(ParsedChunk(section_title=clean_title or None, page_no=page_no, chunk_text=piece))
    return chunks


# ---- txt / md ----------------------------------------------------------------
def _parse_text_file(filepath: str) -> list[tuple[str | None, int | None, str]]:
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    # 마크다운이면 헤더 기준, 순수 txt면 빈 줄(문단) 기준
    if re.search(r"^#{1,6}\s+", content, re.MULTILINE):
        return [(t, None, b) for t, b in _split_by_headers(content)]
    paragraphs = re.split(r"\n\s*\n", content)
    return [(None, None, p) for p in paragraphs if p.strip()]


# ---- pptx (슬라이드 단위) ------------------------------------------------------
def _iter_shape_texts(shape, out: list[str]) -> None:
    """도형에서 텍스트를 재귀적으로 수집한다.
    - 그룹 도형: 내부 도형을 재귀 순회 (안 하면 그룹 안 텍스트가 통째로 누락된다)
    - 표: 셀 텍스트를 행 단위로 추출 (커맨드 표 등 핵심 정보가 표에 있는 경우가 많다)
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 그룹 도형 재귀
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            _iter_shape_texts(sub, out)
        return

    # 표
    if getattr(shape, "has_table", False):
        try:
            rows = []
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                out.append("\n".join(rows))
        except Exception:  # noqa: BLE001
            pass
        return

    # 일반 텍스트 프레임
    if getattr(shape, "has_text_frame", False):
        txt = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
        if txt.strip():
            out.append(txt)


def _parse_pptx(filepath: str, include_speaker_notes: bool = False
                ) -> list[tuple[str | None, int | None, str]]:
    """슬라이드별로 텍스트를 모은다. 슬라이드 제목은 section_title이자 본문 앞머리에도 포함해
    임베딩/전문검색/리랭커가 제목 문맥을 함께 보도록 한다.

    발표자 노트는 기본적으로 제외한다(사내 배포 자료에서 노트는 발표용 메모라
    검색 품질을 떨어뜨리는 경우가 많다). 필요하면 include_speaker_notes=True.
    """
    from pptx import Presentation

    prs = Presentation(filepath)
    sections: list[tuple[str | None, int | None, str]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except Exception:  # noqa: BLE001
            pass

        texts: list[str] = []
        for shape in slide.shapes:
            # 제목 도형은 title로 이미 뽑았으므로 본문에서 중복 수집하지 않는다
            try:
                if title and shape == slide.shapes.title:
                    continue
            except Exception:  # noqa: BLE001
                pass
            _iter_shape_texts(shape, texts)

        if include_speaker_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text
            if note and note.strip():
                texts.append(f"[슬라이드 노트]\n{note.strip()}")

        body = "\n".join(texts).strip()
        if not body and not title:
            continue
        # 제목을 본문 앞에 포함 -> 검색 대상 텍스트에 제목 문맥이 들어간다
        full = f"{title}\n{body}".strip() if title else body
        sections.append((title or f"슬라이드 {idx}", idx, full))
    return sections


# ---- docx / pdf (Docling) ----------------------------------------------------
def _parse_via_docling(filepath: str) -> list[tuple[str | None, int | None, str]]:
    result = _get_converter().convert(filepath)
    markdown = result.document.export_to_markdown()
    return [(t, None, b) for t, b in _split_by_headers(markdown)]


# ---- 공통 헬퍼 ----------------------------------------------------------------
def _split_by_headers(markdown: str) -> list[tuple[str | None, str]]:
    lines = markdown.splitlines()
    sections: list[tuple[str | None, str]] = []
    current_title = None
    current_lines: list[str] = []

    for line in lines:
        if re.match(r"^#{1,6}\s+", line):
            if current_lines:
                sections.append((current_title, "\n".join(current_lines)))
            current_title = re.sub(r"^#{1,6}\s+", "", line).strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_lines:
        sections.append((current_title, "\n".join(current_lines)))
    return sections


def _split_long_text(text: str, max_chars: int) -> list[str]:
    if len(text) <= max_chars:
        return [text]
    paragraphs = text.split("\n\n")
    pieces, buf = [], ""
    for p in paragraphs:
        if len(buf) + len(p) + 2 > max_chars and buf:
            pieces.append(buf)
            buf = p
        else:
            buf = f"{buf}\n\n{p}" if buf else p
    if buf:
        pieces.append(buf)
    # 여전히 너무 긴 조각은 하드 컷
    final = []
    for piece in pieces:
        while len(piece) > max_chars:
            final.append(piece[:max_chars])
            piece = piece[max_chars:]
        if piece:
            final.append(piece)
    return final
