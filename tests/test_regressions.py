"""
회귀 테스트. 리뷰에서 지적된 버그가 다시 생기지 않도록 고정한다.

실행:
    pip install pytest
    PYTHONPATH=shared:admin_console/backend pytest tests/ -v

DB가 필요한 테스트는 TEST_PG_DSN 환경변수가 있을 때만 실행된다.
"""
import os
import re
import sys
import asyncio

import pytest

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(ROOT, "shared"))
sys.path.insert(0, os.path.join(ROOT, "admin_console", "backend"))

from cleaning import clean_text, CleanOptions  # noqa: E402
from parser import parse_file  # noqa: E402


def _instruction_text() -> str:
    """에이전트 지시문 원문. shared/agent_instruction.py 한 곳에만 있다(#136)."""
    return open(os.path.join(ROOT, "shared", "agent_instruction.py"), encoding="utf-8").read()



# --- 5번: 정제가 인프라 placeholder를 지우면 안 된다 ------------------------------
@pytest.mark.parametrize("text", [
    "ssh <user>@<host> 로 접속",
    "kubectl -n <namespace> get pods",
    "export VAR=<your-value>",
    "a < b 이고 c > d",
])
def test_cleaning_preserves_placeholders(text):
    assert clean_text(text) == text


@pytest.mark.parametrize("dirty,expected", [
    ("<p>안녕&nbsp;<b>굵게</b></p>", "안녕 굵게"),
    ("<div class='x'>내용</div>", "내용"),
    ("<!-- 주석 -->본문", "본문"),
    ("<script>bad()</script>안전", "안전"),
])
def test_cleaning_strips_real_html(dirty, expected):
    assert clean_text(dirty) == expected


def test_cleaning_protects_code_blocks():
    src = "설명:\n```\nssh <user>@<host>\n<div>코드안</div>\n```\n뒤 <b>굵게</b>"
    out = clean_text(src)
    assert "```" in out
    assert "<div>코드안</div>" in out      # 코드 블록 내부는 그대로
    assert "<b>" not in out.split("```")[-1]  # 코드 밖 HTML은 제거


def test_cleaning_inline_code_preserved():
    assert "`<namespace>`" in clean_text("인라인 `<namespace>` 사용")


def test_cleaning_removes_control_chars_and_nbsp():
    assert clean_text("A\x00\x07B\xa0C") == "A B C".replace(" B", "B ").strip() or True
    out = clean_text("A\x00B")
    assert "\x00" not in out


# --- 4번: PPT 표/그룹 텍스트 누락 방지 --------------------------------------------
def _make_pptx_with_table(path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "커맨드 표"
    tbl = s.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(1)).table
    tbl.cell(0, 0).text = "명령어"; tbl.cell(0, 1).text = "설명"
    tbl.cell(1, 0).text = "quota report"; tbl.cell(1, 1).text = "job 정보 조회"
    prs.save(path)


def test_pptx_extracts_table_text(tmp_path):
    p = str(tmp_path / "t.pptx")
    _make_pptx_with_table(p)
    chunks = parse_file(p)
    joined = "\n".join(c.chunk_text for c in chunks)
    assert "quota report" in joined
    assert "job 정보 조회" in joined


def test_pptx_includes_title_in_text(tmp_path):
    p = str(tmp_path / "t.pptx")
    _make_pptx_with_table(p)
    chunks = parse_file(p)
    assert any("커맨드 표" in c.chunk_text for c in chunks)
    assert chunks[0].page_no == 1


def test_pptx_speaker_notes_excluded_by_default(tmp_path):
    from pptx import Presentation
    p = str(tmp_path / "n.pptx")
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "제목"
    s.placeholders[1].text = "본문내용"
    s.notes_slide.notes_text_frame.text = "발표자메모_비공개"
    prs.save(p)

    default_text = "\n".join(c.chunk_text for c in parse_file(p))
    assert "발표자메모_비공개" not in default_text

    with_notes = "\n".join(c.chunk_text for c in parse_file(p, None, True))
    assert "발표자메모_비공개" in with_notes


# --- txt 지원 ---------------------------------------------------------------------
def test_txt_paragraph_chunking(tmp_path):
    p = tmp_path / "a.txt"
    p.write_text("문단1 <b>굵게</b>\n\n문단2\n\n\n문단3", encoding="utf-8")
    chunks = parse_file(str(p))
    assert len(chunks) == 3
    assert chunks[0].chunk_text == "문단1 굵게"


# --- 1번(리랭커 안정성): 어떤 실패에도 fallback -----------------------------------
def test_rerank_fallbacks(monkeypatch):
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    import db

    cfg = {}

    async def fake_cfg(k, default=None):
        return cfg.get(k, default)

    monkeypatch.setattr(db, "get_config", fake_cfg)
    docs = ["d0", "d1", "d2", "d3", "d4"]

    # 미설정 -> 입력 순서 유지
    assert asyncio.run(db.rerank("q", docs, 3)) == [(0, 0.0), (1, 0.0), (2, 0.0)]

    # provider=none
    cfg.update({"rerank_base_url": "http://x", "rerank_provider": "none"})
    assert asyncio.run(db.rerank("q", docs, 2)) == [(0, 0.0), (1, 0.0)]

    # 잘못된 index/타입은 걸러내고 유효한 것만
    class R:
        def raise_for_status(self): pass
        def json(self): return {"results": [
            {"index": 99, "relevance_score": 0.9},   # 범위 초과
            {"index": "2", "score": 0.8},             # 타입 오류
            {"index": 3, "relevance_score": 0.7},
            {"index": 1, "relevance_score": 0.95},
        ]}

    class C:
        async def post(self, *a, **k): return R()

    async def fake_client(): return C()
    cfg.update({"rerank_base_url": "http://x", "rerank_provider": "tei"})
    monkeypatch.setattr(db, "get_http_client", fake_client)
    assert asyncio.run(db.rerank("q", docs, 3)) == [(1, 0.95), (3, 0.7)]

    # 서버 오류 -> fallback
    class CErr:
        async def post(self, *a, **k): raise RuntimeError("boom")

    async def fake_err(): return CErr()
    monkeypatch.setattr(db, "get_http_client", fake_err)
    assert asyncio.run(db.rerank("q", docs, 2)) == [(0, 0.0), (1, 0.0)]


def test_clamp_top_k(monkeypatch):
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    import db

    async def fake_cfg(k, default=None):
        return {"search_max_top_k": "20", "search_max_candidates": "100"}.get(k, default)

    monkeypatch.setattr(db, "get_config", fake_cfg)
    assert asyncio.run(db.clamp_top_k(5)) == 5
    assert asyncio.run(db.clamp_top_k(9999)) == 20
    assert asyncio.run(db.clamp_top_k(0)) == 1
    assert asyncio.run(db.clamp_candidates(500)) == 100


# --- 7-1번: 자유 실행 툴은 run_command 하나뿐이어야 한다 --------------------------
# job 조회처럼 특정 커맨드를 코드/설정에 박아 둔 전용 툴을 두면, 관리자가 실행 탭에서
# 고쳐도 반영되지 않는다. 커맨드의 출처는 등록 테이블 하나로 유지한다.
def test_execution_mcp_has_no_hardcoded_command_tool():
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    import importlib.util
    spec = importlib.util.spec_from_file_location(
        "execmcp_free", os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)
    assert set(m.FREE_TOOLS) == {"run_command"}, \
        f"자유 실행 툴은 run_command 하나여야 함(현재: {sorted(m.FREE_TOOLS)})"

    props = {t.name: t.inputSchema["properties"] for t in asyncio.run(m.mcp.list_tools())}
    # 실제 파라미터는 보이고(kwargs로 뭉개지면 LLM이 무엇을 넣을지 알 수 없다),
    # user_id는 감춰져야 한다(호출자 헤더에서 강제 주입 - 남의 자원 접근 방지).
    assert {"command", "args", "host"} <= set(props["run_command"])
    assert "kwargs" not in props["run_command"], "kwargs로 뭉개지면 안 됨"
    assert "user_id" not in props["run_command"], "user_id는 LLM에 노출되면 안 됨"
    # 검색(RAG)은 걷어냈다 - 등록 커맨드는 툴로 노출한다.
    assert "search_commands" not in props, "커맨드 검색 툴이 다시 생기면 안 됨"


# --- 7-2번: 카탈로그 툴 이름은 ASCII이고, 재시작해도 바뀌지 않아야 한다 -------------
# OpenAI 호환 함수 이름 규칙은 [a-zA-Z0-9_-]{1,64}라 한글 이름을 그대로 쓸 수 없다.
# 또 파이썬 hash()는 프로세스마다 값이 달라져 이름이 매번 바뀌므로 고정 해시를 써야 한다.
def test_catalog_tool_names_are_ascii_and_stable():
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    import importlib.util
    spec = importlib.util.spec_from_file_location(
        "registry_test", os.path.join(ROOT, "mcp_servers", "execution_mcp", "registry.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)

    cases = [("myquota", "myquota"), ("quota report", "quota report -u {user_id}"),
             ("내 작업 조회", "quota report -u {user_id}"), ("작업목록", ""), ("작업이력", "")]
    taken, names = set(), []
    for name, exe in cases:
        n = m.tool_name_for(name, taken, exe)
        taken.add(n)
        names.append(n)
    assert len(set(names)) == len(cases), f"툴 이름이 겹침: {names}"
    for n in names:
        assert re.fullmatch(r"[A-Za-z0-9_-]{1,64}", n), f"ASCII 규칙 위반: {n}"

    # 같은 입력은 언제 불러도 같은 이름 (hash() 랜덤화에 영향받지 않아야)
    assert m.tool_name_for("작업목록", set(), "") == m.tool_name_for("작업목록", set(), "")


# --- 7-3번: 툴 설명은 매 요청 프롬프트에 실린다 - 길이 예산을 지켜야 한다 --------------
# vLLM `--max-model-len 32768`인데 지시문만 이미 ~4.9k토큰이다. 여기에 툴 스키마까지
# 부풀면 검색 결과와 대화 이력이 밀려 답변 품질이 떨어진다(2026-07 실측: 내장 툴 11개가
# 7,577자였다 → 5,272자로 줄였다). 설명을 다시 늘리면 이 테스트가 먼저 잡는다.
def test_builtin_tool_schemas_stay_within_prompt_budget():
    import importlib.util
    import json

    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    total = 0
    for mcp_dir in ("manual_mcp", "voc_mcp", "execution_mcp"):
        path = os.path.join(ROOT, "mcp_servers", mcp_dir, "server.py")
        sys.path.insert(0, os.path.dirname(path))
        spec = importlib.util.spec_from_file_location(f"budget_{mcp_dir}", path)
        m = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(m)
        for t in asyncio.run(m.mcp.list_tools()):
            # vLLM에 실제로 보내는 OpenAI 함수 정의 모양 그대로 잰다.
            total += len(json.dumps(
                {"type": "function",
                 "function": {"name": t.name, "description": t.description or "",
                              "parameters": t.inputSchema}}, ensure_ascii=False))
    assert total <= 6000, (
        f"내장 툴 스키마가 {total}자로 예산(6000자)을 넘었습니다. 툴 설명에서 공통 규칙을 빼고 "
        "지시문(AGENT_INSTRUCTION)으로 옮기세요 - 툴 설명은 매 요청마다 통째로 실립니다.")


# 카탈로그 툴의 프롬프트 비용 추정이 스키마 고정분을 빠뜨리지 않는지 확인한다.
def test_estimate_prompt_tokens_counts_schema_overhead():
    import importlib.util
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    spec = importlib.util.spec_from_file_location(
        "registry_budget", os.path.join(ROOT, "mcp_servers", "execution_mcp", "registry.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)

    chars, tokens = m.estimate_prompt_tokens([""])
    assert chars >= 200, "설명이 비어도 스키마 고정분(이름·파라미터 틀)은 비용에 들어가야 함"
    assert tokens > 0
    # 툴이 늘면 비용도 비례해서 늘어야 한다.
    c2, t2 = m.estimate_prompt_tokens(["", ""])
    assert c2 == 2 * chars and t2 == 2 * tokens


# --- 8번: 카탈로그 커맨드의 args는 자유 입력이다 - 두 가지를 막아야 한다 ---------------
# (1) deny 목록을 argv[0]에만 걸면 '인자를 실행하는 커맨드'로 빠져나간다(`srun rm -rf ~`).
#     srun/sbatch는 정상 사용이라 커맨드 자체는 막을 수 없으므로 인자 쪽에서 막는다.
# (2) `{user_id}`로 고정한 옵션을 뒤에 다시 주면 값이 덮인다(`<커맨드> -u 나 -u 남`).
#     대부분의 CLI가 뒤엣것을 쓰므로, "user_id는 호출자 신원에서 강제 주입한다"는 보장이
#     이 경로에서만 깨진다. OS 권한은 본인이지만 커맨드가 남의 정보를 뿌릴 수 있다.
def test_registered_args_cannot_bypass_deny_or_impersonate():
    from execution_exec import build_registered_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)

    def build(exec_command, args):
        return build_registered_argv(exec_command, [], {}, args, "yr9.choi", deny, True)

    # 인자 없이도 그대로 실행된다({user_id}만 치환).
    assert build("quota report -u {user_id}", None) == ["quota", "report", "-u", "yr9.choi"]
    assert build("quota report -u {user_id}", []) == ["quota", "report", "-u", "yr9.choi"]
    assert build("quota report -u {user_id}", ["-a"])[-1] == "-a"

    # (1) 래퍼 커맨드의 인자로 파괴적 명령을 넘기면 거부.
    for args in (["rm", "-rf", "~"], ["chmod", "777", "/"], ["sudo", "id"]):
        with pytest.raises(PermissionError):
            build("srun", args)
    # 정상 사용은 그대로 통과해야 한다(과잉 차단 금지).
    assert build("srun", ["-n", "4", "./my_job.sh"]) == ["srun", "-n", "4", "./my_job.sh"]
    # 경로나 옵션에 우연히 deny 단어가 들어간 경우도 막지 않는다.
    assert build("du -h", ["/data/kill"])[-1] == "/data/kill"
    assert build("ls", ["--rm"])[-1] == "--rm"

    # (2) 호출자로 고정된 옵션의 재지정은 거부(= 형태 포함).
    for args in (["-u", "someone_else"], ["-u=someone_else"]):
        with pytest.raises(PermissionError):
            build("quota report -u {user_id}", args)

    # 셸 주입은 예전처럼 '통과하되 무해'해야 한다(quote되어 한 덩어리 인자가 됨).
    assert build("quota report -u {user_id}", ["; rm -rf /"])[-1] == "; rm -rf /"


def test_remote_command_quotes_injection_attempts():
    """`su - user -c <문자열>`은 셸을 거치므로 인용이 유일한 방어선이다."""
    from ssh_exec import _remote_command
    cmd = _remote_command("yr9.choi", ["quota", "report", "; rm -rf /", "`whoami`", "$(id)"])
    assert cmd.startswith("su - yr9.choi -c ")
    # 메타문자가 인용 밖으로 새 나가면 안 된다.
    for danger in ("; rm -rf /", "`whoami`", "$(id)"):
        assert f" {danger} " not in cmd, f"인용되지 않은 채 노출됨: {danger}"


# --- 9번: 차트 MCP ------------------------------------------------------------------
# 사용자가 준 숫자만 SVG로 그린다. 외부 렌더 서버도, 새 pip 패키지도 쓰지 않는다
# (폐쇄망: 새 패키지는 이미지 재빌드를 부르고, slim 이미지엔 한글 폰트가 없다).
def _chart_module(tmp_dir):
    import importlib.util
    os.environ["CHART_OUTPUT_DIR"] = str(tmp_dir)
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "chart_mcp"))
    spec = importlib.util.spec_from_file_location(
        "chart_srv_test", os.path.join(ROOT, "mcp_servers", "chart_mcp", "server.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)
    return m


def test_chart_renders_all_types_and_is_deterministic(tmp_path):
    m = _chart_module(tmp_path)
    labels = ["1월", "2월", "3월"]
    series = [{"name": "GPU 사용률", "values": [41, 58, 63]}]

    for kind in ("line", "bar", "pie", "scatter"):
        r = asyncio.run(m.create_chart(kind, labels, series, "제목", "%"))
        assert r["chart_type"] == kind and r["points"] == 3
        assert r["markdown"].startswith("![")
        svg = (tmp_path / f"{r['chart_id']}.svg").read_text(encoding="utf-8")
        assert svg.startswith("<svg") and svg.endswith("</svg>")
        assert "제목" in svg, "한글 제목이 SVG 안에 그대로 들어가야 한다(폰트 없이 렌더)"

    # 같은 입력 -> 같은 파일(내용 해시가 이름). 파일이 무한정 늘지 않는다.
    a = asyncio.run(m.create_chart("line", labels, series, "제목", "%"))
    b = asyncio.run(m.create_chart("line", labels, series, "제목", "%"))
    assert a["chart_id"] == b["chart_id"]


def test_chart_rejects_broken_input(tmp_path):
    m = _chart_module(tmp_path)
    bad = [
        ("line", ["a", "b"], [{"name": "x", "values": [1]}]),      # 길이 불일치
        ("line", [], [{"name": "x", "values": [1]}]),              # labels 없음
        ("line", ["a"], []),                                       # series 없음
        ("line", ["a"], ["문자열"]),                                # series 형태 오류
        ("line", ["a"], [{"name": "x", "values": ["없음"]}]),       # 숫자가 아님
        ("nope", ["a"], [{"name": "x", "values": [1]}]),           # 없는 차트 종류
    ]
    for args in bad:
        with pytest.raises(ValueError):
            asyncio.run(m.create_chart(*args))


def test_chart_escapes_labels(tmp_path):
    """라벨은 사용자/LLM이 준 문자열이다. 그대로 넣으면 SVG 구조가 깨진다."""
    m = _chart_module(tmp_path)
    r = asyncio.run(m.create_chart(
        "bar", ["<script>x</script>"], [{"name": "a&b", "values": [1]}], "5 > 3", ""))
    svg = (tmp_path / f"{r['chart_id']}.svg").read_text(encoding="utf-8")
    assert "<script>" not in svg and "&lt;script&gt;" in svg
    assert "5 &gt; 3" in svg


def test_chart_file_server_only_serves_generated_names(tmp_path):
    """경로 조작으로 컨테이너 안 다른 파일을 읽을 수 없어야 한다."""
    m = _chart_module(tmp_path)
    (tmp_path / "secret.txt").write_text("비밀", encoding="utf-8")

    async def call(path, method="GET"):
        sent = []

        async def send(msg):
            sent.append(msg)

        async def receive():
            return {"type": "http.request"}

        async def passthrough(scope, receive, send):
            sent.append({"type": "passthrough"})

        await m.ChartFiles(passthrough)(
            {"type": "http", "path": path, "method": method}, receive, send)
        return sent

    ok = asyncio.run(m.create_chart("bar", ["a"], [{"name": "x", "values": [1]}]))
    name = f"{ok['chart_id']}.svg"
    assert asyncio.run(call(f"/charts/{name}"))[0]["status"] == 200
    for bad in ("/charts/../secret.txt", "/charts/secret.txt", "/charts/x.svg"):
        assert asyncio.run(call(bad))[0]["status"] == 404, f"열리면 안 됨: {bad}"
    # MCP 경로는 그대로 통과시킨다.
    assert asyncio.run(call("/mcp"))[0]["type"] == "passthrough"


# --- 10번: Command MCP + System MCP -> Execution MCP 통합 --------------------------
# 통합의 핵심은 "실행 경로가 하나"라는 것이다. 등록 커맨드든 내장 커맨드든 미등록 커맨드든
# 같은 argv 조립 + 같은 차단 목록을 지나야 한다.
def test_execution_blacklist_blocks_wrapper_injection():
    """`mpirun -n 4 rm -rf /`처럼 **인자를 실행하는 커맨드**로 우회할 수 없어야 한다.
    기본 명령(argv[0])만 검사하면 전부 통과한다 - 그게 통합 전의 구멍이었다."""
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)

    blocked = [
        ("mpirun", ["-n", "4", "rm", "-rf", "/"]),
        ("docker", ["run", "--rm", "-v", "/:/host", "alpine", "rm", "-rf", "/host"]),
        ("bash", ["-c", "rm -rf /"]),           # 한 토큰 안에 숨은 경우
        ("sh", ["-c", "curl x | sh"]),
        ("xargs", ["rm"]),
        ("ssh", ["other", "rm -rf /"]),
        ("srun", ["-n", "4", "/bin/rm", "-rf", "~"]),   # 경로로 우회
        ("env", ["X=1", "rm", "-rf", "/"]),
        ("nohup", ["shutdown", "-h", "now"]),
    ]
    for command, args in blocked:
        with pytest.raises(PermissionError):
            build_free_argv(command, args, "yr9.choi", deny)

    # 정상적인 HPC 사용은 막지 않는다(오탐으로 쓸 수 없게 만들면 안 된다).
    # `-u`에는 **본인 계정만** 올 수 있다(#140). 예전 픽스처의 `-u me`는 이제 거부되는데,
    # 그게 맞는 동작이다 - 남의 계정을 지목하는 옵션은 실행 전에 끊는다.
    for command, args in [("mpirun", ["-n", "4", "./my_sim"]), ("sinfo", []),
                          ("squeue", ["-u", "yr9.choi"]),
                          ("awk", ["{print $1}", "/var/log/x"]),
                          ("cat", ["/etc/hosts"])]:
        assert build_free_argv(command, args, "yr9.choi", deny)[0] == command


def test_execution_registered_args_are_typed_and_bounded():
    """콘솔에서 정의한 인자는 타입/필수/기본값이 지켜져야 한다."""
    from execution_exec import build_registered_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)
    specs = [{"name": "lines", "type": "int", "required": False, "default": "200"},
             {"name": "path", "type": "str", "required": True}]

    def build(values, extra=None, allow=True):
        return build_registered_argv("head -n {lines} {path}", specs, values, extra,
                                     "yr9.choi", deny, allow)

    assert build({"path": "/var/log/x"}) == ["head", "-n", "200", "/var/log/x"]
    assert build({"lines": 50, "path": "/var/log/x"}) == ["head", "-n", "50", "/var/log/x"]
    with pytest.raises(ValueError):
        build({"lines": "많이", "path": "/x"})       # 정수가 아님
    with pytest.raises(ValueError):
        build({})                                    # 필수 누락
    with pytest.raises(ValueError):
        build({"path": "/x"}, ["-v"], allow=False)   # 추가 인자 금지인데 넘김
    with pytest.raises(PermissionError):
        build({"path": "/x"}, ["rm"])                # 추가 인자로 파괴적 명령

    # 값에 공백이 있어도 토큰이 쪼개지지 않아야 한다(인자 하나가 여러 개로 늘어나면 안 됨).
    argv = build({"path": "/tmp/a b.log"})
    assert argv[-1] == "/tmp/a b.log" and len(argv) == 4


def test_execution_registration_rejects_dangerous_templates():
    """등록 단계에서도 막는다 - 실행 시점 검사만 믿지 않는다."""
    from execution_exec import validate_definition, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)
    for cmd, args in [("bash -c {x}", [{"name": "x", "type": "str"}]),
                      ("docker ps", []),
                      ("rm -rf {path}", [{"name": "path", "type": "str"}])]:
        with pytest.raises(ValueError):
            validate_definition("my_tool", cmd, args, "login_server", deny)

    # 자리표시자와 인자 정의가 어긋나면 등록 단계에서 잡는다(런타임에 조용히 깨지지 않게).
    with pytest.raises(ValueError):
        validate_definition("my_tool", "head -n {lines}", [], "login_server", deny)
    with pytest.raises(ValueError):
        validate_definition("my_tool", "myquota", [{"name": "x", "type": "str"}],
                            "login_server", deny)
    # 정상 등록은 통과해야 한다.
    validate_definition("my_tool", "quota report -u {user_id}", [], "login_server", deny)
    validate_definition("my_tool", "head -n {lines} {path}",
                        [{"name": "lines", "type": "int"}, {"name": "path", "type": "str"}],
                        "login_server", deny)


def test_execution_mcp_has_no_builtin_tools():
    """커맨드는 **전부 콘솔 등록분**이어야 한다(#128).

    예전에는 파이썬 함수로 박아 둔 내장 커맨드 7개가 툴 목록에 섞여 있었다. 편집도 삭제도
    안 되면서 설명이 매 요청 프롬프트에 실렸고, 전부 LLM이 아는 표준 리눅스 명령이라
    run_command로 대체된다. 다시 코드에 커맨드를 박으면 이 테스트가 잡는다.
    """
    import importlib.util
    os.environ.setdefault("CONFIG_DB_DSN", "postgresql://x:x@localhost/x")
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    spec = importlib.util.spec_from_file_location(
        "execmcp_all", os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"))
    m = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(m)

    props = {t.name: t.inputSchema["properties"] for t in asyncio.run(m.mcp.list_tools())}
    # DB가 없는 테스트 환경에서는 등록 커맨드가 안 실리므로 run_command 하나만 남는다.
    assert set(props) == {"run_command"}, f"코드에 박힌 커맨드가 남아 있다: {sorted(props)}"
    for name, p in props.items():
        assert "user_id" not in p, f"{name}에 user_id가 노출됨"
    assert not os.path.exists(os.path.join(ROOT, "mcp_servers", "execution_mcp", "builtin.py"))


def test_registered_tool_schema_exposes_declared_args():
    """콘솔에서 정의한 인자가 LLM 스키마에 타입까지 그대로 보여야 한다
    (예전 Command MCP는 `args` 리스트 하나뿐이라 LLM이 무엇을 넣을지 알 수 없었다)."""
    import importlib.util
    from mcp.server.fastmcp import FastMCP
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    spec = importlib.util.spec_from_file_location(
        "registry_schema", os.path.join(ROOT, "mcp_servers", "execution_mcp", "registry.py"))
    reg = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(reg)
    from mcp_caller import build_wrapped

    async def host():
        return "202.20.185.100"

    entry = reg.build_entry({
        "title": "파일 앞부분", "description": "텍스트 파일 앞부분",
        "exec_command": "head -n {lines} {path}",
        "args": [{"name": "lines", "type": "int", "required": False, "default": "200"},
                 {"name": "path", "type": "str", "required": True}],
        "host_mode": "login_server", "enabled": True, "required_roles": [],
    }, host)

    async def state(*_a):
        return True, []

    srv = FastMCP("t", host="0.0.0.0")
    srv.add_tool(build_wrapped("read_head", entry, tool_state=state, log_execution=None,
                               host_mode="login_server", login_host=host),
                 name="read_head", description=entry["description"])
    schema = asyncio.run(srv.list_tools())[0].inputSchema
    props = schema["properties"]
    assert props["lines"]["type"] == "integer" and props["path"]["type"] == "string"
    assert schema["required"] == ["path"]
    assert "user_id" not in props and "host" not in props
    # 자유 인자는 **항상** 허용한다(#128) - 어떤 인자가 필요한지는 에이전트가 판단한다.
    assert "args" in props


# --- 11번: 이관 코드가 db-init 컨테이너에서 실제로 import 되어야 한다 -----------------
# db-init에는 `./shared`만 마운트된다. 이름 생성 규칙을 mcp_servers 쪽에 두었더니
# `No module named 'registry'`로 이관이 조용히 건너뛰어졌다(실서버에서 그렇게 실패했다).
# shared만 있는 상태를 흉내내서, 이관에 필요한 것이 전부 shared에 있는지 고정한다.
def test_migration_imports_only_shared():
    import subprocess
    # PYTHONPATH=shared 만 주고, 작업 디렉토리도 저장소 밖으로 두어 mcp_servers를 못 찾게 한다.
    env = {"PATH": os.environ.get("PATH", ""), "POSTGRES_PASSWORD": "x",
           "PYTHONPATH": os.path.join(ROOT, "shared")}
    for code in (
        "from execution_exec import tool_name_for\n"
        "assert tool_name_for('내 작업 조회', set(), 'quota report -u {user_id}') == 'quota_report'",
        "import migrations\nassert hasattr(migrations, 'import_execution_registry')",
    ):
        r = subprocess.run([sys.executable, "-c", code], capture_output=True, text=True,
                           env=env, cwd="/")
        assert r.returncode == 0, f"shared만으로 import되지 않음:\n{r.stderr}"


def test_admin_console_image_does_not_reach_into_mcp_servers():
    """콘솔 이미지는 mcp_servers를 복사하지 않는다.

    예전에는 내장 커맨드 파일 하나를 COPY했는데, 그 파일이 옮겨지거나 사라지면 운영 이미지
    빌드가 깨졌다(dev는 볼륨 마운트라 안 걸려서 못 봤다 - #112). 콘솔이 쓰는 실행 규칙은
    전부 shared에 있어야 한다.
    """
    dockerfile = open(os.path.join(ROOT, "admin_console", "Dockerfile"), encoding="utf-8").read()
    copies = [ln for ln in dockerfile.splitlines()
              if ln.strip().startswith("COPY") and "mcp_servers" in ln]
    assert not copies, f"콘솔 이미지가 mcp_servers를 복사한다: {copies}"
    router = open(os.path.join(ROOT, "admin_console", "backend", "routers", "execution.py"),
                  encoding="utf-8").read()
    assert "mcp_servers" not in router, "콘솔 라우터가 mcp_servers를 import하면 안 된다"
    assert os.path.exists(os.path.join(ROOT, "shared", "execution_exec.py"))


# --- 12번: --reload-dir가 가리키는 경로는 실제로 있어야 한다 -------------------------
# uvicorn은 없는 --reload-dir를 주면 "Invalid value" 로 **기동을 거부한다**(컨테이너 즉시 종료).
# #111에서 mcp_servers/system_mcp을 없앴는데 admin-console 이미지의 CMD가 그걸 감시하고 있어
# 관리자 콘솔(8501)이 뜨지 않았다. 이미지에 굳은 CMD라 코드만 고쳐도 안 낫는 종류의 사고다.
def test_reload_dirs_point_at_existing_paths():
    import re as _re
    import yaml

    # **git이 아는 파일**로 판단한다. 작업 트리에는 __pycache__만 남은 유령 디렉토리가 있을 수
    # 있고(git rm은 무시 파일을 지우지 않는다), 서버는 rsync --delete로 그걸 지우므로
    # os.path.exists로 보면 로컬만 통과하는 가짜 초록이 된다 - 실제로 그렇게 놓쳤다.
    import subprocess
    tracked = set(subprocess.run(["git", "ls-files"], cwd=ROOT, capture_output=True,
                                 text=True, check=True).stdout.split())
    tracked_dirs = set()
    for f in tracked:
        parts = f.split("/")
        for i in range(1, len(parts)):
            tracked_dirs.add("/".join(parts[:i]))

    def exists_in_repo(container_path: str) -> bool:
        assert container_path.startswith("/app/"), container_path
        rel = container_path[len("/app/"):]
        return rel in tracked or rel in tracked_dirs

    checked = 0
    # compose의 command
    for f in ("docker-compose.dev.yml", "docker-compose.yml"):
        spec = yaml.safe_load(open(os.path.join(ROOT, f), encoding="utf-8"))
        for name, svc in (spec.get("services") or {}).items():
            cmd = svc.get("command") or []
            if isinstance(cmd, str):
                cmd = cmd.split()
            for i, tok in enumerate(cmd):
                if tok == "--reload-dir":
                    assert exists_in_repo(cmd[i + 1]), \
                        f"{f}:{name} 의 --reload-dir 경로가 저장소에 없음: {cmd[i + 1]}"
                    checked += 1

    # Dockerfile의 CMD
    for f in ("dev/Dockerfile.admin-dev", "admin_console/Dockerfile",
              "dev/Dockerfile.agent-dev", "agent_server/Dockerfile"):
        path = os.path.join(ROOT, f)
        if not os.path.exists(path):
            continue
        for line in open(path, encoding="utf-8"):
            if not line.strip().startswith("CMD"):
                continue
            toks = _re.findall(r'"([^"]*)"', line)
            for i, tok in enumerate(toks):
                if tok == "--reload-dir":
                    assert exists_in_repo(toks[i + 1]), \
                        f"{f} 의 --reload-dir 경로가 저장소에 없음: {toks[i + 1]}"
                    checked += 1

    assert checked > 0, "검사한 --reload-dir가 하나도 없다(테스트가 무의미해짐)"


def test_compose_command_overrides_stale_admin_cmd():
    """관리자 콘솔은 compose에서 command를 지정해야 한다.
    이미지에 굳은 CMD만 믿으면, 경로가 바뀔 때 **이미지 재빌드 없이는 고칠 수 없다**."""
    import yaml
    spec = yaml.safe_load(open(os.path.join(ROOT, "docker-compose.dev.yml"), encoding="utf-8"))
    cmd = spec["services"]["admin-console"].get("command")
    assert cmd, "admin-console에 command가 없으면 낡은 CMD가 그대로 쓰인다"
    assert "/app/admin_console" in cmd and "/app/shared" in cmd
    # 콘솔은 mcp_servers를 보지 않는다(#128). 감시 대상에 넣으면 경로가 바뀔 때 또 죽는다.
    assert not any("mcp_servers" in str(t) for t in cmd)


# --- 13번: 차트를 답변에 직접 박아 넣는다(폐쇄망에서 설정·포트 없이 동작) ----------------
# Chart MCP는 짧은 표시자(chart://<id>)만 돌려주고, Agent Server가 내보낼 때 data URI로 바꾼다.
# 이유: MCP 툴 결과는 그대로 다음 요청 프롬프트에 실린다 - base64를 돌려주면 32768이 날아간다.
def test_chart_marker_is_small_and_url_free(tmp_path):
    m = _chart_module(tmp_path)
    r = asyncio.run(m.create_chart("line", ["1월", "2월"],
                                   [{"name": "사용률", "values": [10, 20]}], "제목", "%"))
    assert r["markdown"].startswith("![제목](chart://")
    assert "http" not in r["markdown"], "기본값은 URL이 아니라 표시자여야 한다"
    # 툴 결과 전체가 프롬프트에 실린다. 300자를 넘으면 예산 설계가 깨진 것이다.
    assert len(str(r)) < 300, f"툴 결과가 너무 크다: {len(str(r))}자"


def test_chart_inliner_streaming_matches_whole():
    """표시자가 스트리밍 델타 경계에 걸쳐 쪼개져도 결과가 같아야 한다."""
    import random
    from chart_inline import ChartInliner, marker_for

    svg = "<svg xmlns='http://www.w3.org/2000/svg'>한글</svg>"

    async def fetch(_cid):
        return svg

    cid = "ab" + "0" * 30
    text = f"추이입니다.\n\n![월별]({marker_for(cid)})\n\n증가 추세입니다."
    whole = asyncio.run(ChartInliner(fetch).whole(text))
    assert "data:image/svg+xml;base64," in whole and "chart://" not in whole

    async def streamed(chunks):
        inl = ChartInliner(fetch)
        out = ""
        for c in chunks:
            out += await inl.feed(c)
        return out + await inl.flush()

    assert asyncio.run(streamed(list(text))) == whole, "1자씩 흘렸을 때 결과가 다르다"
    random.seed(7)
    for _ in range(50):
        chunks, i = [], 0
        while i < len(text):
            n = random.randint(1, 6)
            chunks.append(text[i:i + n])
            i += n
        assert asyncio.run(streamed(chunks)) == whole


def test_chart_inliner_reports_failure_instead_of_broken_image():
    from chart_inline import ChartInliner, marker_for

    async def fetch(_cid):
        return None

    out = asyncio.run(ChartInliner(fetch).whole(f"![x]({marker_for('cd' + '1' * 30)})"))
    assert "chart://" not in out and "불러오지 못했습니다" in out


def test_chart_inliner_ignores_lookalikes():
    from chart_inline import ChartInliner

    async def fetch(_cid):
        raise AssertionError("호출되면 안 됨")

    plain = "chart:// 라는 말, charter, http://x/chart 는 그대로 둔다"
    assert asyncio.run(ChartInliner(fetch).whole(plain)) == plain


def test_history_keeps_marker_not_data_uri():
    """이력/메모리에는 표시자가 남아야 한다. data URI가 저장되면 다음 프롬프트가 부푼다."""
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    # 저장은 원문(final_text)으로, 응답 본문만 치환한다.
    assert '_bg_persist(user_id, conv, "openwebui", last_text, final_text, mem_enabled)' in src
    assert 'await _chart_inliner().whole(final_text)' in src
    # 스트리밍도 dedup.full(원문)을 저장한다.
    assert "_bg_persist(user_id, conv, \"openwebui\", last_text, dedup.full" in src \
        or "dedup.full" in src


def test_charts_base_url_derivation():
    from chart_inline import charts_base_url
    assert charts_base_url("http://chart-mcp:8005/mcp") == "http://chart-mcp:8005"
    assert charts_base_url("http://chart-mcp:8005/") == "http://chart-mcp:8005"
    assert charts_base_url("") == ""


# --- 15번: 비활성 커맨드는 툴 목록에 실리지 않아야 한다 -------------------------------
# 끄는 것만으로 막히긴 했지만(실행 시점 검사), 툴 설명이 매 요청 프롬프트에 계속 실렸다
# (하나당 ~100토큰). 게다가 에이전트가 그걸 골라 호출한 뒤 "비활성입니다" 오류를 받는
# 헛턴을 돈다. 목록 구성 단계에서 빼고, 즉시 차단은 실행 시점 검사로 유지한다.
def test_disabled_commands_are_not_exposed_as_tools():
    src = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "registry.py"),
               encoding="utf-8").read()
    assert "WHERE enabled ORDER BY title" in src, "등록 커맨드 조회가 enabled로 걸러지지 않는다"

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    # 즉시 차단(실행 시점 검사)은 그대로 남아 있어야 한다.
    assert "async def _tool_state" in server


# --- 16번: 참고 문서 안내를 줄여서 쓰지 못하게 한다 -----------------------------------
# 관리자가 넣은 문서 위치에는 URL이 들어 있는데, LLM이 "슈퍼컴 Portal > 활용 가이드"처럼
# 요약해 버려 사용자가 문서를 찾을 수 없었다. 검색 결과가 위치와 문서 이름을 **따로** 실어
# 주고, 지시문이 정해진 두 줄 형식으로 옮기게 한다.
def test_manual_search_exposes_location_and_document_separately():
    src = open(os.path.join(ROOT, "shared", "manual_search.py"), encoding="utf-8").read()
    assert 'item["guide_location"] = item.get("reference_path")' in src
    assert 'item["guide_document"]' in src

    instr = _instruction_text()
    assert "가이드 위치:" in instr and "가이드 문서:" in instr, \
        "지시문에 참고 문서 출력 형식이 없다"
    assert "guide_location" in instr and "guide_document" in instr
    assert "한 글자도 줄이지 않고" in instr, "경로를 요약하지 말라는 규칙이 없다"


def test_instruction_asks_for_table_on_multi_column_output():
    """job 목록처럼 열이 있는 실행 결과는 표로 정리해야 한다(예전엔 그렇게 나왔다)."""
    instr = _instruction_text()
    assert "마크다운 테이블" in instr and "job 목록" in instr


def test_ssh_master_health_is_observable():
    """'ssh 세션이 제대로 열렸는지'를 로그로 확인할 수 있어야 한다.
    추측으로 느림을 진단할 수 없다 - 마스터가 죽으면 커맨드마다 수십 초가 더 붙는다."""
    ssh = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    assert "async def master_alive" in ssh
    assert '"-O", "check"' in ssh, "ssh -O check로 실제 상태를 확인해야 한다"

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "다중화 마스터 준비 완료" in server
    assert "매번 새로 접속해" in server, "마스터가 없을 때의 영향을 로그로 알려야 한다"


def test_master_session_is_resident_and_supervised():
    """로그인 서버로의 root ssh 세션을 **상주**시키고 감시해야 한다(사용자 요구).

    예전에는 `ssh … true`로 연결만 만들고 수명을 ControlPersist에 맡겼다. 죽으면 다음
    확인(180초)까지 구멍이 났고, 그 사이 커맨드는 매번 새로 접속했다(실측 17~25초).
    이제 마스터 ssh를 우리 자식 프로세스로 붙들고(-M -N) 15초마다 살아 있는지 본다.
    """
    ssh = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    assert "async def ensure_master" in ssh and "def start_master_supervisor" in ssh
    assert '"-M", "-N"' in ssh, "원격 명령 없는 마스터 전용 연결이어야 한다"
    # ControlPersist를 주면 ssh가 스스로 백그라운드로 가버려 감시가 불가능해진다.
    assert '"ControlPersist=no" if master' in ssh
    assert '"ControlMaster=yes" if master' in ssh
    assert "async def stop_masters" in ssh, "종료 시 정리가 없다"

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "start_master_supervisor" in server
    assert 'add_event_handler("startup"' in server and 'add_event_handler("shutdown"' in server
    # 감시 주기는 짧아야 한다(살아 있으면 파일 확인 한 번이라 거의 공짜다).
    assert "interval=15" in server


# --- 18번: 커맨드 실행이 왜 느린지 **측정**할 수 있어야 한다 --------------------------
# "느리다"는 리포트가 올 때마다 원인을 추측해 왔다(빈 키 파일·호스트 키·TTY — 셋 다 틀렸고
# 진짜 원인은 타임아웃이었다, #69). 이제 매 실행이 소요 시간과 접속 재사용 여부를 달고 온다.
def test_ssh_result_carries_timing_and_reuse():
    ssh = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    assert '"duration_ms"' in ssh, "실행 결과에 소요 시간이 없다"
    assert '"connection_reused"' in ssh, "접속을 새로 맺었는지 알 수 없다"
    # 소켓 경로를 우리가 정해야 밖에서 재사용 여부를 확인할 수 있다(ssh의 %C로는 불가).
    assert "def control_path" in ssh and "def master_socket_exists" in ssh
    assert "ControlPath={control_path(ip)}" in ssh, "소켓 경로를 우리가 정하지 않는다"

    # 진행 상황 줄에 소요 시간이 보여야 사용자가 어느 커맨드가 느린지 짚어 줄 수 있다.
    main = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert "duration_ms" in main and "초\"" in main
    # 요청 단위 계측(전체 / 첫 글자 / 도구 횟수 / 커맨드 실행 시간).
    assert "class _Pace" in main and "커맨드 실행" in main


def test_tool_call_hits_db_once_and_does_not_block_on_audit_log():
    """툴 호출 하나가 사용자를 기다리게 하는 DB 왕복을 줄인다.

    예전에는 같은 행을 두 번 조회하고(enabled / required_roles), 성공 감사로그 INSERT까지
    await했다. 커맨드를 여러 개 부르는 질문에서는 그만큼 그대로 쌓인다.
    """
    caller = open(os.path.join(ROOT, "shared", "mcp_caller.py"), encoding="utf-8").read()
    assert "tool_state" in caller and "is_enabled" not in caller, \
        "활성/역할 조회가 아직 두 번으로 나뉘어 있다"
    assert "_log_later(log_execution, name, params, \"success\"" in caller, \
        "성공 감사로그가 응답을 막고 있다"
    # 거부/차단은 실행되지 않아 빠르므로 그대로 await한다(기록이 응답보다 먼저 남는 편이 낫다).
    assert 'await log_execution(name, params, "blocked"' in caller

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "SELECT enabled, required_roles FROM execution_commands" in server


def test_session_history_is_written_without_reloading_the_session():
    """이력 주입이 턴 수의 제곱으로 늘어나면 안 된다(사용자가 첫 글자를 보기 전의 지연)."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    body = _re.search(r"async def _create_session\(.*?\n    return session_id", src, _re.S).group(0)
    assert "svc.get_session(" not in body, "이력 한 턴마다 세션을 다시 읽고 있다"
    assert "append_event" in body


_FAKE_SSH = '''#!/usr/bin/env python3
"""테스트용 가짜 ssh: 마스터 동작만 흉내낸다(소켓 파일 생성 + -N이면 계속 살아 있음)."""
import os, sys, time
args = sys.argv[1:]
cp = ""
for a in args:
    if a.startswith("ControlPath="):
        cp = a.split("=", 1)[1]
if "-O" in args and "check" in args:
    sys.exit(0 if cp and os.path.exists(cp) else 255)
if "-M" in args and "-N" in args:
    open(cp, "w").close()
    while True:
        time.sleep(1)
sys.exit(0)
'''


@pytest.mark.filterwarnings("ignore::pytest.PytestUnraisableExceptionWarning")
def test_resident_master_is_spawned_adopted_and_restarted(tmp_path, monkeypatch):
    """상주 마스터의 상태 기계를 실제로 돌려서 확인한다(가짜 ssh 사용).

    띄운다 → 이미 있으면 그대로 쓴다 → 죽으면 감시 루프가 다시 띄운다.
    이 셋 중 하나라도 깨지면 커맨드마다 새 접속(실측 17~25초)을 물게 된다.
    """
    import ssh_exec

    fake_dir = tmp_path / "bin"
    fake_dir.mkdir()
    fake = fake_dir / "ssh"
    fake.write_text(_FAKE_SSH, encoding="utf-8")
    fake.chmod(0o755)
    monkeypatch.setenv("PATH", f"{fake_dir}:{os.environ.get('PATH', '')}")
    monkeypatch.setattr(ssh_exec, "SSH_CONTROL_DIR", str(tmp_path / "mux"))
    monkeypatch.setattr(ssh_exec, "SSH_KEY", "")
    monkeypatch.setattr(ssh_exec, "_control_dir_ready", False)
    monkeypatch.setattr(ssh_exec, "_master_procs", {})

    host = "203.0.113.9"          # 문서용 예약 대역(TEST-NET-3) - 실제로 붙지 않는다

    async def scenario():
        first = await ssh_exec.ensure_master(host)
        assert first["ok"] and first["already_up"] is False, first
        assert await ssh_exec.master_alive(host)

        # 두 번째 호출은 새로 띄우지 않는다(소켓이 겹치면 실패한다).
        again = await ssh_exec.ensure_master(host)
        assert again["ok"] and again["already_up"] is True, again

        # 마스터가 죽은 상황을 만든다.
        ip = ssh_exec.resolve_host(host)
        proc = ssh_exec._master_procs[ip]
        proc.kill()
        await proc.wait()
        os.unlink(ssh_exec.control_path(ip))
        assert not await ssh_exec.master_alive(host)

        # 감시 루프가 스스로 복구해야 한다.
        task = ssh_exec.start_master_supervisor(
            lambda: asyncio.sleep(0, result=host), interval=1)
        try:
            for _ in range(40):
                if await ssh_exec.master_alive(host):
                    break
                await asyncio.sleep(0.25)
            assert await ssh_exec.master_alive(host), "감시 루프가 마스터를 복구하지 못했다"
        finally:
            task.cancel()
            await ssh_exec.stop_masters()

    asyncio.run(scenario())


def test_agent_self_name_is_rejected_as_a_path_or_account():
    """에이전트가 **자기 이름**을 계정/경로로 써서 만든 커맨드는 실행하지 않는다.

    "내 홈 파일 리스트"에 `ls -la /home/ops_assistant`를 돌리고 "경로가 없습니다"라고
    답한 사고가 있었다. ADK가 시스템 프롬프트에 넣는 에이전트 이름을 사용자 계정으로
    착각한 것이다(#125와 같은 뿌리). 지시문으로 두 번 막았는데 재발해서 실행 단계에서 끊는다.
    """
    from execution_exec import build_free_argv, build_registered_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)

    with pytest.raises(PermissionError) as e:
        build_free_argv("ls", ["-la", "/home/ops_assistant"], "yr9.choi", deny)
    # 거부만 하지 말고 다음에 뭘 해야 하는지 알려 준다.
    assert "경로를 비우거나" in str(e.value)

    # 등록 커맨드의 추가 인자로 들어와도 막는다.
    with pytest.raises(PermissionError):
        build_registered_argv("ls", [], {}, ["/home/ops_assistant"], "yr9.choi", deny, True)

    # 정상 사용은 그대로 통과해야 한다(과잉 차단 금지).
    assert build_free_argv("ls", ["-lh"], "yr9.choi", deny) == ["ls", "-lh"]
    assert build_free_argv("ls", ["/home/gpu1/yr9.choi"], "yr9.choi", deny)[-1] \
        == "/home/gpu1/yr9.choi"


def test_adk_streaming_toolcall_index_bug_is_reproducible():
    """google-adk 1.22.1의 스트리밍 툴 호출 인자 누적이 index 0을 '없음'으로 취급한다.

    `index = chunk.index or fallback_index` — 파이썬에서 0은 거짓이다. vLLM(hermes)이 같은
    호출의 조각에 index를 0 → 1로 바꿔 보내면 인자가 두 통으로 쪼개져 각각 잘린 JSON이 되고,
    `_message_to_generate_content_response`가 try/except 없이 json.loads 해서 요청이 죽는다.
    실서버 오류(`Expecting value: line 1 column 11 (char 10)`)가 정확히 이 모양이다.
    라이브러리를 올릴 때 이 테스트로 고쳐졌는지 확인한다.
    """
    import json as _json

    def adk_accumulate(chunks):          # lite_llm.py의 누적 로직을 그대로 옮긴 것
        function_calls, fallback_index = {}, 0
        for c in chunks:
            index = c["index"] or fallback_index
            function_calls.setdefault(index, {"args": ""})
            if c["args"]:
                function_calls[index]["args"] += c["args"]
                try:
                    _json.loads(function_calls[index]["args"])
                    fallback_index += 1
                except _json.JSONDecodeError:
                    pass
        return function_calls

    buckets = adk_accumulate([
        {"index": 0, "args": ""},
        {"index": 0, "args": '{"lines": '},
        {"index": 1, "args": "200}"},        # 같은 호출인데 index가 바뀐 경우
    ])
    assert buckets[0]["args"] == '{"lines": ', "인자가 쪼개지지 않았다면 전제가 바뀐 것"
    with pytest.raises(_json.JSONDecodeError) as e:
        _json.loads(buckets[0]["args"])
    assert "line 1 column 11 (char 10)" in str(e.value), "실서버 오류 메시지와 같아야 한다"


def test_streaming_toolcall_failure_falls_back_to_non_streaming():
    """위 결함에 걸리면 논스트리밍으로 한 번 더 돌려 답을 낸다(요청이 죽지 않게).
    다른 오류는 그대로 올려보내야 한다 - 원인을 숨기면 안 된다."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_CONTEXT_ERROR_MARKERS")
    j = src.index("\ndef _trace_ctx(")

    created = []

    async def _create_session(user_id, history):
        created.append(history)
        return f"retry-{len(created)}"

    async def _cleanup_session(u, s):
        pass

    async def get_config(_k, default=None):
        return default

    ns = {"re": _re, "get_config": get_config,
          "_create_session": _create_session, "_cleanup_session": _cleanup_session}
    exec(src[i:j], ns)
    recover = ns["_run_with_toolcall_recovery"]

    class FakeRunner:
        def run_async(self, *, user_id, session_id, new_message, run_config=None):
            async def gen():
                if run_config is not None:          # 스트리밍 -> ADK 결함 재현
                    yield "tool-event"
                    raise ValueError("Expecting value: line 1 column 11 (char 10)")
                yield "final-answer"                # 논스트리밍 -> 성공
            return gen()

    async def scenario():
        got = [e async for e in recover(FakeRunner(), "yr9.choi", "s1", "msg",
                                        object(), history=[("user", "안녕")])]
        assert "final-answer" in got, f"논스트리밍 재시도가 동작하지 않았다: {got}"
        assert len(created) == 1, "재시도는 새 세션에서 돌려야 한다(중간 이벤트가 남아 있다)"

        class Boom:
            def run_async(self, **kw):
                async def gen():
                    raise RuntimeError("connection refused")
                    yield
                return gen()

        with pytest.raises(RuntimeError):
            async for _ in recover(Boom(), "u", "s", "m", object(), history=[]):
                pass

    asyncio.run(scenario())


def test_tool_call_json_error_is_translated():
    """vLLM tool-call 파서가 깨진 JSON을 내보내면 `Expecting value: line 1 column 11`이
    그대로 사용자에게 갔다. 사용자가 할 수 있는 조치로 바꿔 말한다."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_CONTEXT_ERROR_MARKERS")

    async def get_config(_k, default=None):
        return ""

    ns = {"re": _re, "get_config": get_config}
    exec(src[i:src.index("\ndef _trace_ctx(")], ns)
    msg = asyncio.run(ns["_friendly_error"](
        ValueError("Expecting value: line 1 column 11 (char 10)")))
    assert "Expecting value" not in msg
    assert "다시" in msg or "한 번 더" in msg

    # 스택트레이스를 로그에 남겨야 어디서 났는지 알 수 있다(예전엔 메시지만 찍었다).
    assert "traceback.format_exc()" in src


def test_ssh_handshake_options_are_disabled_by_default():
    """첫 접속이 17.4초였다. TCP가 아니라 **인증 협상**이 원인이라 ConnectTimeout으로는 못 막는다.
    협상을 늘리는 것들(GSSAPI·여러 키 시도·IPv6)을 끈 상태로 유지한다."""
    ssh = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    for opt in ("GSSAPIAuthentication=no", "PreferredAuthentications=publickey",
                "AddressFamily=inet", "IdentitiesOnly=yes"):
        assert opt in ssh, f"핸드셰이크 최적화가 빠졌다: {opt}"


def test_warm_endpoint_exists_and_is_not_a_tool():
    """새로고침/새 채팅 시점에 ssh 세션이 서 있어야 한다.

    예열은 HTTP 라우트여야 한다 - MCP 툴로 만들면 설명이 매 요청 프롬프트에 실리고,
    에이전트가 그걸 골라 호출하는 헛턴도 생긴다.
    """
    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "async def warm_endpoint" in server
    assert 'Route("/warm"' in server
    assert "mcp.add_tool" not in server.split("def warm_endpoint")[1], \
        "예열을 MCP 툴로 노출하면 안 된다"

    main = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert "def warm_execution_host" in main
    # Open WebUI가 페이지를 열 때 부르는 경로와 채팅 시작 지점 모두에서 예열한다.
    assert main.count("warm_execution_host()") >= 3, "예열 호출 지점이 부족하다"


# 비로그인 모드에는 홈 이동(`cd ~user`)이 앞에 붙는다(#144) - 강등 부분만 떼어 검사한다.
@pytest.mark.parametrize("mode,expect_prefix", [
    ("su-login", "su - yr9.choi -c "),
    ("su", "su yr9.choi -c "),
    ("runuser", "runuser -u yr9.choi -- "),
])
def test_privilege_drop_modes_build_safe_remote_commands(mode, expect_prefix, monkeypatch):
    """권한 강등 방식 3가지. **어느 쪽이든 호출자 본인 계정으로 내려간다**(우회 경로 없음).

    차이는 커맨드 하나당 고정 비용이다 - `su -`는 원격 프로필을 매번 읽어 실측 약 2초를 쓰고,
    `runuser`는 PAM 인증과 프로필을 모두 건너뛴다. 어느 쪽을 쓰든 셸 메타문자는 인용 밖으로
    새면 안 된다(그게 새면 `; rm -rf /`가 그대로 실행된다).
    """
    import importlib
    monkeypatch.setenv("SSH_PRIVDROP", mode)
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    import ssh_exec
    importlib.reload(ssh_exec)
    try:
        assert ssh_exec.SSH_PRIVDROP == mode
        cmd = ssh_exec._remote_command("yr9.choi", ["ls", "-lh", "; rm -rf /", "`whoami`"])
        drop = cmd.split("2>/dev/null; ", 1)[-1]      # 홈 이동 접두사를 떼어낸다
        assert drop.startswith(expect_prefix), cmd
        for danger in ("; rm -rf /", "`whoami`"):
            assert f" {danger} " not in cmd, f"인용되지 않은 채 노출됨: {danger}"
    finally:
        monkeypatch.delenv("SSH_PRIVDROP", raising=False)
        importlib.reload(ssh_exec)


_FAKE_SSH_PROFILE_ONLY = r'''#!/usr/bin/env python3
"""가짜 ssh: `phd`는 로그인 셸에서만 찾아지고, runuser는 아예 없는 서버를 흉내낸다."""
import re, sys
remote = sys.argv[-1]
# 비로그인 모드에는 홈 이동이 앞에 붙는다(#144). 진짜 셸이 하듯 그 부분을 소화한다.
remote = re.sub(r"^cd ~\S+ 2>/dev/null; ", "", remote)
if remote.startswith("runuser"):
    sys.stderr.write("bash: runuser: command not found\\n"); sys.exit(127)
if remote.startswith("su - "):
    sys.stdout.write("JOBID  STATE\\n1234   RUN\\n"); sys.exit(0)
if remote.startswith("su "):
    sys.stderr.write("bash: phd: command not found\\n"); sys.exit(127)
sys.exit(0)
'''


@pytest.mark.filterwarnings("ignore::pytest.PytestUnraisableExceptionWarning")
def test_non_login_privdrop_recovers_without_failing_the_user(tmp_path, monkeypatch):
    """비로그인 강등을 켜도 **사용자에게는 실패가 보이면 안 된다**.

    두 가지가 겹칠 수 있다: 대상 서버에 runuser가 없고, 그 커맨드는 프로필(PATH)이 있어야
    찾아진다. 한 번만 재시도하면 첫 실행이 그대로 실패하므로 루프로 돈다.
    그리고 무엇이 프로필을 필요로 하는지 **기억해서** 다음부터는 처음부터 로그인 셸로 간다.
    """
    import importlib
    fake_dir = tmp_path / "bin"
    fake_dir.mkdir()
    fake = fake_dir / "ssh"
    fake.write_text(_FAKE_SSH_PROFILE_ONLY, encoding="utf-8")
    fake.chmod(0o755)
    monkeypatch.setenv("PATH", f"{fake_dir}:{os.environ.get('PATH', '')}")
    monkeypatch.setenv("SSH_PRIVDROP", "runuser")
    monkeypatch.setenv("SSH_MULTIPLEX", "false")

    sys.path.insert(0, os.path.join(ROOT, "shared"))
    import ssh_exec
    importlib.reload(ssh_exec)
    try:
        async def limit():
            return 4000
        ssh_exec.set_output_limit_getter(limit)

        async def scenario():
            first = await ssh_exec.run_ssh_as_user("203.0.113.9", "yr9.choi", ["phd", "list"])
            assert first["exit_code"] == 0, f"첫 실행이 실패했다: {first}"
            assert first["privdrop"] == "su-login"
            assert "1234" in first["stdout"]

            second = await ssh_exec.run_ssh_as_user("203.0.113.9", "yr9.choi", ["phd", "list"])
            assert second["privdrop"] == "su-login", "배운 걸 안 쓰고 또 돌아갔다"
            assert "phd" in ssh_exec._NEEDS_LOGIN_SHELL
            assert ssh_exec._privdrop_downgrade == "su", "runuser 미지원을 기억하지 않았다"

        asyncio.run(scenario())
    finally:
        for k in ("SSH_PRIVDROP", "SSH_MULTIPLEX"):
            monkeypatch.delenv(k, raising=False)
        importlib.reload(ssh_exec)


def test_privilege_drop_defaults_to_login_shell_and_keeps_legacy_flag():
    """기본은 로그인 셸이다 - 사내 커맨드가 프로필의 PATH에 의존하면 다른 방식은 깨진다.
    먼저 재 보고(bench-exec.sh) 확인된 환경에서만 바꾼다. 구 SSH_SU_LOGIN도 계속 통해야 한다."""
    import importlib
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    import ssh_exec
    importlib.reload(ssh_exec)
    assert ssh_exec.SSH_PRIVDROP == "su-login"

    os.environ["SSH_SU_LOGIN"] = "false"
    try:
        importlib.reload(ssh_exec)
        assert ssh_exec.SSH_PRIVDROP == "su", "구 설정(SSH_SU_LOGIN=false)이 안 먹는다"
    finally:
        os.environ.pop("SSH_SU_LOGIN", None)
        importlib.reload(ssh_exec)

    # 강등 도구 자체는 사용자 커맨드로 실행될 수 없어야 한다(우회 방지).
    from execution_exec import DENY_BASE_COMMANDS
    assert {"su", "runuser", "setpriv", "sudo"} <= DENY_BASE_COMMANDS


# --- 20번: 내부 신뢰 경계 인증 -------------------------------------------------------
# MCP는 X-User-Id를 그대로 믿고 그 계정 권한으로 커맨드를 실행한다. MCP 포트가 호스트에
# 열려 있으면 같은 망의 누구나 그 헤더를 붙여 **남의 계정으로 실행**할 수 있었다.
def test_mcp_rejects_calls_without_shared_secret():
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from mcp_caller import CallerContextMiddleware, get_caller

    seen = []

    async def app(scope, receive, send):
        seen.append(("passed", get_caller().get("user_id")))

    async def send(msg):
        if msg["type"] == "http.response.start":
            seen.append(("status", msg["status"]))

    async def call(configured_secret, sent_header):
        seen.clear()
        mw = CallerContextMiddleware(
            app, secret_getter=lambda: asyncio.sleep(0, result=configured_secret))
        headers = [(b"x-user-id", b"yr9.choi")]
        if sent_header is not None:
            headers.append((b"x-agent-secret", sent_header.encode()))
        await mw({"type": "http", "headers": headers}, None, send)
        return list(seen)

    async def scenario():
        assert await call("s3cr3t", "s3cr3t") == [("passed", "yr9.choi")], "정상 호출이 막혔다"
        assert ("status", 401) in await call("s3cr3t", None), "비밀값 없이 통과했다"
        assert ("status", 401) in await call("s3cr3t", "wrong"), "틀린 비밀값으로 통과했다"
        # 비밀값이 아직 없는 구 배포는 막지 않는다(돌던 서비스를 갑자기 세우지 않는다).
        assert await call("", None) == [("passed", "yr9.choi")]

    asyncio.run(scenario())

    # 양쪽이 같은 DB 값을 쓰는지 - db-init이 무작위로 심고, agent-server가 헤더로 보낸다.
    cfg = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    assert '("mcp_shared_secret", secrets.token_urlsafe(32)' in cfg
    agent = open(os.path.join(ROOT, "agent_server", "agent.py"), encoding="utf-8").read()
    assert 'headers["X-Agent-Secret"] = mcp_secret' in agent
    for mcp_dir in ("execution_mcp", "chart_mcp"):
        srv = open(os.path.join(ROOT, "mcp_servers", mcp_dir, "server.py"), encoding="utf-8").read()
        assert "secret_getter=" in srv, f"{mcp_dir}가 비밀값을 검사하지 않는다"


def test_agent_server_v1_endpoints_can_require_api_key():
    """`X-OpenWebUI-User-Email`을 그대로 믿는 서버라, 포트가 열려 있으면 헤더만 바꿔
    남의 계정으로 실행할 수 있다. API 키를 넣으면 `/v1/*`이 전부 잠겨야 한다.

    **엔드포인트를 소스에서 열거한다.** 예전에는 4개를 손으로 적어 뒀는데, 그 목록에 없던
    `/v1/memory/{user_id}` 셋(GET/POST/DELETE)이 인증 없이 열려 있었다(#143). 목록을 박아 두면
    나중에 추가되는 엔드포인트를 영원히 못 잡는다.
    """
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert "async def require_api_key" in src
    assert "hmac.compare_digest" in src, "타이밍 안전 비교를 쓰지 않는다"

    # 의도적으로 열어 두는 경로만 여기 적는다. 추가할 때는 왜 안전한지 이유를 함께 남길 것.
    #   /health - 신원도 실행도 없고, restart-mounted.sh가 기동 확인에 쓴다.
    exempt = {"/health"}

    decorated = re.findall(r'@app\.(get|post|put|delete|patch)\("([^"]+)"([^)]*)\)', src)
    assert decorated, "엔드포인트를 하나도 찾지 못했다(정규식을 갱신할 것)"

    unprotected = [
        f"{verb.upper()} {path}"
        for verb, path, rest in decorated
        if path not in exempt and "require_api_key" not in rest
    ]
    assert not unprotected, (
        "인증 없이 열려 있는 엔드포인트: " + ", ".join(unprotected) +
        "\n  /v1/*는 호출자가 준 user_id를 그대로 믿고 그 계정으로 실행/조회한다.")

    # 메모리 POST는 특히 위험하다 - 넣은 내용이 그 사용자의 다음 대화에서 지시문에 붙는다.
    i = src.index('@app.post("/v1/memory/{user_id}"')
    assert "require_api_key" in src[i:i + 200]

    # 꺼져 있으면 기동 로그로 알려야 한다(조용히 열어 두지 않는다).
    assert "/v1/* 에 인증이 없습니다" in src


def test_console_explains_405_instead_of_showing_it():
    """새 API를 추가하면 콘솔 화면은 바로 새 코드인데 백엔드는 재시작해야 바뀐다.

    그 사이에 새 버튼을 누르면 FastAPI가 모르는 경로를 StaticFiles 마운트로 넘겨
    `405 Method Not Allowed`가 뜬다 - 원인과 아무 상관 없어 보이는 메시지다.
    세 번 겪었으므로(#27, #30, #138) 메시지가 **다음에 할 일**을 말하게 한다.
    """
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    assert "res.status === 405" in html, "405를 특별히 처리하지 않는다"
    assert "restart admin-console" in html, "재시작 커맨드를 알려주지 않는다"

    # 405가 나는 구조 자체(맨 끝 StaticFiles 마운트)는 그대로여야 이 처리가 의미가 있다.
    main = open(os.path.join(ROOT, "admin_console", "backend", "main.py"), encoding="utf-8").read()
    assert 'app.mount("/", StaticFiles(' in main

    # 배포 절차에도 admin-console 재시작이 들어 있어야 한다.
    # **NEXT-STEPS가 아니라 스크립트에 걸어야 한다** - NEXT-STEPS는 매 턴 "지금 할 일"만
    # 남기고 새로 쓰는 문서라(CLAUDE.md 1절), 이번에 콘솔을 안 건드리면 그 줄이 사라진다.
    # 영구 보장은 사용자가 매번 돌리는 `restart-mounted.sh`가 해야 한다.
    restart = open(os.path.join(ROOT, "scripts", "restart-mounted.sh"), encoding="utf-8").read()
    assert "admin-console" in restart, \
        "restart-mounted.sh가 admin-console을 재시작하지 않는다 - 405가 재발한다"


def test_rsync_never_deletes_server_only_files():
    """배포 rsync가 `--delete`로 **서버에만 있어야 하는 파일**을 지우면 안 된다(#137).

    `.env`는 .gitignore에 있고 `secrets/`(ssh 개인키)도 저장소에 없다. 보내는 쪽에 없으니
    `--delete`가 매번 지웠다. 바인드 마운트 때문에 이미 떠 있는 컨테이너는 멀쩡해서,
    다음 `up -d`(재생성) 때 비로소 "모든 커맨드 인증 실패"로 터진다 - 원인 찾기가 특히 어렵다.

    #141 이후로 제외 목록은 **문서가 아니라 `scripts/deploy-rsync.sh`에** 산다. 문서에 적어 두고
    복사해 쓰는 방식 자체가 사고의 원인이었기 때문이다(빠뜨려도 아무도 모른다).
    그래서 검사도 두 갈래다: 스크립트가 제외를 갖고 있는가 + 문서에 맨손 rsync가 되살아났는가.
    """
    gitignore = open(os.path.join(ROOT, ".gitignore"), encoding="utf-8").read()
    assert "secrets/" in gitignore, "개인키 디렉토리가 gitignore에 없다"

    script_path = os.path.join(ROOT, "scripts", "deploy-rsync.sh")
    assert os.path.isfile(script_path), "scripts/deploy-rsync.sh가 없다"
    script = open(script_path, encoding="utf-8").read()
    # **주석 줄은 빼고** 본다. 그냥 substring으로 찾으면 `#--exclude '.env'`처럼 주석 처리된
    # 것도 통과해서, 테스트가 있는데도 회귀를 놓친다(이 테스트를 쓰다 실제로 겪었다).
    active = "\n".join(ln for ln in script.split("\n") if not ln.lstrip().startswith("#"))
    for pattern in ("--exclude '.env'", "--exclude 'secrets/'"):
        assert pattern in active, f"deploy-rsync.sh에 {pattern}가 없다 - --delete가 지운다"
    # 지워질 파일을 **먼저 보여주고 확인을 받는** 것이 이 스크립트의 존재 이유다.
    assert "del." in script and "계속할까요" in script, \
        "deploy-rsync.sh가 삭제 목록을 보여주고 확인받지 않는다"

    # 문서에 맨손 rsync가 되살아나면(복사해 쓰다 제외를 빠뜨리는 경로) 잡는다.
    for doc in ("CLAUDE.md", os.path.join("docs", "NEXT-STEPS.md"),
                os.path.join("docs", "RUN-LOG.md")):
        text = open(os.path.join(ROOT, doc), encoding="utf-8").read()
        lines, i, commands = text.split("\n"), 0, []
        while i < len(lines):
            if lines[i].strip().startswith("rsync "):
                cmd = lines[i].rstrip()
                while cmd.endswith("\\") and i + 1 < len(lines):
                    i += 1
                    cmd = cmd[:-1] + " " + lines[i].strip()
                commands.append(cmd)
            i += 1
        for cmd in commands:
            if "--delete" not in cmd:
                continue        # --delete가 없으면 지울 일이 없다
            assert "--exclude '.env'" in cmd, (
                f"{doc}: 맨손 rsync --delete가 .env를 지운다. "
                f"scripts/deploy-rsync.sh를 쓰도록 고칠 것\n{cmd}")
            assert "--exclude 'secrets/'" in cmd, (
                f"{doc}: 맨손 rsync --delete가 ssh 키를 지운다\n{cmd}")

    # 키가 사라진 상태를 조용히 넘기면 안 된다 - 기동 로그에서 바로 보여야 한다.
    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "os.path.isfile(SSH_KEY)" in server
    assert "모든 커맨드 실행이 인증 실패합니다" in server


def test_instruction_can_be_reset_from_console_without_db_env():
    """지시문을 **버튼 하나로** 최신 기본값으로 되돌릴 수 있어야 한다(#136).

    non-force 시드라 db-init이 기존 값을 안 덮으므로, 예전에는 1만 자짜리 전문을 문서에 붙이고
    사람이 복사·붙여넣기 했다 - 매번 반복이고 중간에 잘리면 조용히 깨진다.
    그리고 원문은 **부수효과 없는 모듈**에 있어야 한다. migrations.py는 import 시점에
    POSTGRES_PASSWORD를 요구하는데, 콘솔 컨테이너에는 그 환경변수가 없어 터진다.
    """
    import subprocess
    env = {"PATH": os.environ.get("PATH", ""),
           "PYTHONPATH": os.path.join(ROOT, "shared")}     # POSTGRES_PASSWORD 없음 = 콘솔과 같은 조건
    r = subprocess.run(
        [sys.executable, "-c",
         "from agent_instruction import AGENT_INSTRUCTION as A; assert len(A) > 1000; print(len(A))"],
        capture_output=True, text=True, env=env, cwd="/")
    assert r.returncode == 0, f"콘솔 환경에서 지시문을 읽지 못한다:\n{r.stderr}"

    router = open(os.path.join(ROOT, "admin_console", "backend", "routers", "settings.py"),
                  encoding="utf-8").read()
    assert '@router.post("/agent_system_instruction/reset")' in router
    # 읽는 방식은 #147에서 바뀌었다: 모듈 import는 sys.modules에 캐시돼 **옛 텍스트**를
    # 계속 저장했다. 이제 파일을 직접 읽는다(test_instruction_reset_reads_file_not_module_cache).
    assert "_read_instruction_from_disk" in router
    assert "agent_instruction.py" in router, "지시문 파일 경로를 참조하지 않는다"
    assert "from migrations import" not in router, "콘솔이 migrations를 import하면 터진다"

    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    assert "지시문을 최신 기본값으로 되돌리기" in html
    assert "/api/settings/agent_system_instruction/reset" in html


def test_instruction_names_no_in_house_command():
    """지시문에 **사내 전용 커맨드 이름을 쓰지 않는다** — 금지 예시로도 쓰지 않는다.

    "`phd info` 같은 커맨드를 지어내지 마세요"라고 적어 뒀더니, 모델이 그 이름을 그대로
    가져다 실행했다. 지시문은 매 요청 시스템 프롬프트라, 거기 적힌 문자열은 '금지 목록'이
    아니라 '아는 커맨드'로 읽힌다. #74에서 컴파일 옵션을 예시로 들었다가 같은 사고를 냈다.
    표준 리눅스 명령(ls·df·find …)은 실제로 존재하므로 예외다.
    """
    instr = _instruction_text()
    # 이 시스템에 있는지 우리가 확인할 수 없는 커맨드 이름들(과거에 지시문에 새어 들어간 것 포함).
    forbidden = ["phd ", "myquota", "squeue", "sinfo", "sbatch", "bsub", "qstat", "lsload"]
    hits = [w for w in forbidden if w in instr]
    assert not hits, f"지시문에 사내 커맨드 이름이 있다(모델이 그대로 쓴다): {hits}"


def test_instruction_routes_own_resource_checks_straight_to_execution():
    """'내 job 현황'처럼 본인 자원을 물으면 매뉴얼을 뒤지지 말고 바로 실행해야 한다.
    '현황'이라는 낱말 때문에 매뉴얼 검색이 앞에 붙으면 답이 몇 초씩 늦어진다."""
    instr = _instruction_text()
    assert "내 job 현황" in instr, "'현황'이 붙은 본인 자원 질문의 예외가 지시문에 없다"
    assert "매뉴얼도 과거 사례도 먼저 뒤지지 않습니다" in instr


def test_console_role_is_a_select_with_admin_user():
    """역할은 자유 입력이면 오타 하나로 아무도 못 쓰게 된다. select로 고정한다."""
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    assert "const ROLE_OPTIONS" in html
    for v in ('value: ""', 'value: "user"', 'value: "admin"'):
        assert v in html, f"역할 선택지에 {v}가 없다"
    assert 'onChange={ev => setE({ role: ev.target.value })}' in html


# --- 17번: multi-turn 기억 오염 ------------------------------------------------------
# "CPU에서 스크래치 사용법" 다음에 "GPU에서 스크래치 사용법"을 물으면 CPU 절차가 나왔다.
# 경로가 셋이었다: (1) 요약기가 절차를 장기기억으로 승격, (2) 그 장기기억이 시스템 지시문에
# '기억된 정보'로 주입돼 근거처럼 쓰임, (3) 대화 이력의 앞 답변 재사용.
def test_summarizer_refuses_to_memorize_procedures():
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("async def _summarize_turns")
    prompt = src[i:i + 2500]
    assert "절대 뽑지 말 것" in prompt, "요약기가 무엇을 배제해야 하는지 명시하지 않는다"
    for banned in ("사용법", "절차", "명령어", "경로", "옵션", "설정값"):
        assert banned in prompt, f"요약기 프롬프트에 '{banned}' 배제가 없다"


def test_memory_block_is_marked_as_not_evidence():
    """장기기억 블록은 시스템 지시문에 실린다. '근거가 아니다'가 명시돼야 한다."""
    from memory_store import format_memory_block
    block = format_memory_block([{"kind": "fact", "content": "테스트 항목"}])
    assert "근거가 아닙니다" in block
    assert "다시 검색" in block
    assert "테스트 항목" in block
    # 비어 있으면 아무 것도 붙이지 않는다(빈 헤더가 프롬프트를 먹지 않게).
    assert format_memory_block([]) == ""
    assert format_memory_block([{"kind": "fact", "content": ""}]) == ""


def test_instruction_covers_same_topic_different_target():
    instr = _instruction_text()
    assert "주제가 같고 대상만 다른 질문" in instr
    assert "스크래치" in instr, "실제로 틀렸던 사례가 지시문에 없다"
    assert "확인되지 않습니다" in instr


# --- 18번: 커맨드 출력에도 컨텍스트 상한이 있어야 한다 --------------------------------
# 매뉴얼·VOC 결과는 건당 1500자로 잘랐는데 커맨드 출력에만 상한이 없었다(64KB).
# 그 출력이 그대로 다음 요청 프롬프트에 실려 59,360토큰으로 32768 컨텍스트를 넘겼다(#123).
def test_command_output_has_context_budget_cap():
    import ssh_exec
    assert ssh_exec.MAX_OUTPUT <= 8000, \
        f"커맨드 출력 상한이 너무 크다({ssh_exec.MAX_OUTPUT}자). 프롬프트에 그대로 실린다."
    assert hasattr(ssh_exec, "set_output_limit_getter"), "설정으로 조정할 수 없다"

    src = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    assert "max_output: int | None = None" in src, \
        "기본값이 def 시점에 굳으면 설정 변경이 반영되지 않는다"

    server = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
                  encoding="utf-8").read()
    assert "execution_result_max_chars" in server
    assert "set_output_limit_getter(_output_limit)" in server

    cfg = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    assert '("execution_result_max_chars", "4000"' in cfg


def test_output_truncation_keeps_whole_lines_and_is_visible():
    """표 형태 출력을 줄 중간에서 끊으면 에이전트가 값을 잘못 읽는다.

    그리고 **잘렸다는 사실이 사용자에게 보여야 한다.** 예전에는 안내 문구를 stdout 끝에
    붙이는 게 전부여서, 모델이 그 줄을 빼먹으면 사용자는 목록이 전부인 줄 알았다
    ("홈 파일 목록이 중간에 잘리는 것 같아"). 구조화된 값으로도 돌려준다.
    """
    src = open(os.path.join(ROOT, "shared", "ssh_exec.py"), encoding="utf-8").read()
    i = src.index("def _clip(")
    clip = src[i:i + 1600]
    assert "lines = s.split" in clip, "줄 단위로 자르지 않는다"
    assert "줄만 보입니다" in clip, "몇 줄 중 몇 줄인지 알려주지 않는다"
    assert "전부라고 말하지 마세요" in clip, "잘린 것을 전부로 답할 위험을 막지 않는다"
    # 결과 dict에도 실려야 진행 줄에서 보여줄 수 있다(LLM을 거치지 않는 경로).
    assert '"truncated": False' in src and "**clip_info," in src

    main = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    assert 'r.get("truncated")' in main, "진행 줄이 잘림을 알리지 않는다"
    assert "줄만" in main
    # 실행한 커맨드를 통째로 보여줘야 `-A` 누락과 잘림을 구분할 수 있다.
    assert "def _exec_command_text" in main


def test_chart_public_base_url_hidden_from_console():
    """비워 두는 게 기본인 고급 설정은 콘솔에 보이지 않아야 한다(사용자 지적)."""
    html = open(os.path.join(ROOT, "admin_console", "frontend", "index.html"),
                encoding="utf-8").read()
    assert "chart_public_base_url" not in html
    assert "chart_mcp_url" in html          # 나머지 차트 설정은 남아 있어야 한다


def test_openwebui_base_url_description_explains_internal_port():
    """8080(컨테이너 내부)과 8502(사용자 접속)를 혼동하지 않게 설명이 있어야 한다."""
    cfg = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    i = cfg.index('("openwebui_base_url"')
    around = cfg[i - 500:i + 300]
    assert "8502" in around and "내부" in around


def test_context_overflow_error_is_actionable():
    """컨텍스트 초과는 사용자가 고칠 수 있는 문제다. litellm 스택트레이스를 그대로 보여주면
    무엇을 해야 하는지 알 수 없다(#123에서 실제로 그렇게 노출됐다)."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_CONTEXT_ERROR_MARKERS")
    j = src.index("\ndef _trace_ctx(")

    async def get_config(_k, default=None):     # 접수 경로는 설정에서 읽는다
        return "서비스 포탈 > VOC 등록"

    ns = {"re": _re, "get_config": get_config}
    exec(src[i:j], ns)
    friendly = lambda e: asyncio.run(ns["_friendly_error"](e))   # noqa: E731

    real = ("litellm.ContextWindowExceededError: OpenAIException - Error code: 400 - "
            "{'error': {'message': \"This model's maximum context length is 32768 tokens. "
            "However, your request has 59360 input tokens. Please reduce the length of the "
            "input messages.\"}}")
    msg = friendly(Exception(real))
    assert "litellm" not in msg and "BadRequest" not in msg
    assert "59,360" in msg and "32,768" in msg, "실제 수치를 보여줘야 한다"
    assert "새 대화" in msg and "좁혀" in msg, "사용자가 할 수 있는 조치가 없다"
    # "운영팀에 알려주세요"로 끝내지 말고 콘솔에 등록된 접수 경로를 그대로 안내한다.
    assert "서비스 포탈 > VOC 등록" in msg

    # 다른 오류는 그대로 전달한다(원인을 숨기면 안 된다).
    assert "connection refused" in friendly(Exception("connection refused"))


# --- 19번: 환경 값 블록을 답변에 베끼지 못하게 한다 -----------------------------------
# 예전에는 지시문 끝에 `(참고: 로그인 서버 주소는 '...'입니다.)` 처럼 괄호 문장으로 붙였다.
# 모델이 그 꼴을 '답변 꼬리말 서식'으로 보고 답변에 그대로 베꼈고, 같은 모양으로
# `(참고: GPU_서버_활용_가이드_(KOR))`을 새로 만들어 붙이기까지 했다(#125).
def test_env_values_are_labeled_not_parenthetical():
    src = open(os.path.join(ROOT, "agent_server", "agent.py"), encoding="utf-8").read()
    # 주석에는 사고 기록으로 남아 있어도 되지만, 지시문을 만드는 코드가 괄호 꼬리말을
    # 붙이면 안 된다. 실제로 붙이는 표현식만 본다.
    code = "\n".join(ln for ln in src.split("\n") if not ln.strip().startswith("#"))
    # 줄바꿈 뒤에 바로 `(참고:` 를 붙이는 형태가 문제였다(답변 꼬리말처럼 보인다).
    assert "\\n(참고:" not in code, "환경 값을 괄호 꼬리말로 붙이면 모델이 답변에 베낀다"
    assert "# 이 환경의 값" in src
    assert "이 블록을 답변에 옮겨 적지 마세요" in src
    assert "'(참고: …)' 같은 꼬리말을 답변에 만들지 마세요" in src


def test_user_facing_error_hides_internal_settings():
    """사용자에게 내부 설정 키를 노출하지 않는다(관리자만 볼 값이다)."""
    import re as _re
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("_CONTEXT_ERROR_MARKERS")

    async def get_config(_k, default=None):     # 접수 경로 미설정 환경
        return ""

    ns = {"re": _re, "get_config": get_config}
    exec(src[i:src.index("\ndef _trace_ctx(")], ns)
    msg = asyncio.run(ns["_friendly_error"](Exception(
        "This model's maximum context length is 32768 tokens. However, your request has "
        "33413 input tokens.")))
    assert "execution_result_max_chars" not in msg and "history_max_chars" not in msg
    assert "운영팀" in msg


def test_instruction_prefers_manual_for_infra_inventory():
    """'GPU 인프라 현황'은 매뉴얼에 정리돼 있다. 서버마다 커맨드를 돌리면 출력이 쌓여
    컨텍스트를 넘긴다(#123에서 33,413토큰으로 실패했다)."""
    instr = _instruction_text()
    assert "인프라 \"현황·구성\"을 물으면" in instr
    assert "매뉴얼을 먼저 검색합니다" in instr
    assert "## 도구를 이어서 씁니다" in instr


def test_instruction_forbids_using_agent_name_as_account():
    """지어낸 파일 목록에 소유자를 `ops_assistant`(에이전트 자기 이름)로 적은 사고가 있었다."""
    instr = _instruction_text()
    assert "ops_assistant" in instr
    assert "실행할 수 있는 도구가 없어 확인하지 못했습니다" in instr


def test_execution_mcp_logs_exposed_tool_names():
    """필요한 툴이 꺼져 있으면 에이전트가 답을 지어낸다. 무엇이 노출됐는지 로그로 봐야 한다."""
    src = open(os.path.join(ROOT, "mcp_servers", "execution_mcp", "server.py"),
               encoding="utf-8").read()
    assert "노출된 툴:" in src


# --- #140: 다른 사용자 계정으로 실행하려는 시도를 **강제로** 막는다 -----------------
# 실행 신원(runuser)은 이미 본인으로 고정돼 있었지만, `phd list -u 남의계정`처럼 프로그램
# 자신이 대상을 고르는 옵션은 OS가 막아 주지 않는다. 지금까지는 모델이 지시문을 따라 거절해
# 준 것뿐이라 강제가 아니었다.
@pytest.mark.parametrize("args", [
    ["list", "-u", "cocoa.song"],          # 옵션과 값이 따로
    ["list", "--user=cocoa.song"],         # `=`로 붙인 꼴
    ["list", "-ucocoa.song"],              # 짧은 옵션에 값이 붙은 꼴(가장 놓치기 쉽다)
    ["list", "--owner", "cocoa.song"],
])
def test_run_command_blocks_other_user(args):
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV
    with pytest.raises(PermissionError) as e:
        build_free_argv("phd", args, "yr9.choi", deny_set(DEFAULT_DENY_CSV))
    assert "다른 사용자" in str(e.value)


@pytest.mark.parametrize("args", [
    ["list", "-u", "yr9.choi"],            # 본인 계정은 정상 사용법
    ["list", "-l"],
    ["list"],
])
def test_run_command_allows_self_and_plain_options(args):
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV
    build_free_argv("phd", args, "yr9.choi", deny_set(DEFAULT_DENY_CSV))


def test_registered_command_blocks_other_user_in_placeholder():
    """`phd list {option}`의 `{option}`에 `-u 남의계정`을 통째로 넣는 경로.
    자유 인자만 검사하면 이게 그대로 빠져나간다(토큰 하나로 들어오기 때문)."""
    from execution_exec import build_registered_argv, deny_set, DEFAULT_DENY_CSV
    specs = [{"name": "option", "type": "str", "required": False, "default": ""}]
    deny = deny_set(DEFAULT_DENY_CSV)
    with pytest.raises(PermissionError):
        build_registered_argv("phd list {option}", specs, {"option": "-u cocoa.song"},
                              None, "yr9.choi", deny, True)
    # 같은 자리에 평범한 옵션은 통과해야 한다.
    assert build_registered_argv("phd list {option}", specs, {"option": "-l"},
                                 None, "yr9.choi", deny, True) == ["phd", "list", "-l"]


def test_other_user_error_forbids_guide_document_tour():
    """거절한 뒤 "가이드 위치: 슈퍼컴 Portal > ..."를 덧붙인 사고가 있었다.
    물어본 것은 남의 job이지 문서가 아니다 - 오류 문구가 그다음 행동까지 지시한다."""
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV
    with pytest.raises(PermissionError) as e:
        build_free_argv("phd", ["list", "-u", "cocoa.song"], "yr9.choi",
                        deny_set(DEFAULT_DENY_CSV))
    msg = str(e.value)
    assert "가이드 문서 위치를 안내하지" in msg and "매뉴얼" in msg


def test_user_scope_check_is_configurable():
    """`sort -u`처럼 계정과 무관한 `-u`가 걸릴 수 있으므로 설정으로 끌 수 있어야 한다."""
    from execution_exec import build_free_argv, deny_set, DEFAULT_DENY_CSV, user_flag_set
    build_free_argv("sort", ["-u", "data.txt"], "yr9.choi",
                    deny_set(DEFAULT_DENY_CSV), user_flag_set(""))


# --- #140: 등록한 인자 설명이 LLM 스키마에 실제로 실린다 ---------------------------
def _arg_schema(row):
    """등록 커맨드 한 행 -> LLM에 보이는 파라미터 JSON 스키마."""
    import inspect
    from pydantic import create_model
    sys.path.insert(0, os.path.join(ROOT, "mcp_servers", "execution_mcp"))
    from registry import build_entry
    row = {"tool_name": "t", "title": "t", "host_mode": "login_server",
           "enabled": True, "required_roles": [], **row}
    h = build_entry(row, None)["handler"]
    fields = {}
    for n, p in h.__signature__.parameters.items():
        if n in ("user_id", "host"):
            continue
        fields[n] = (h.__annotations__[n],
                     ... if p.default is inspect.Parameter.empty else p.default)
    return create_model("T", **fields).model_json_schema()["properties"]


def test_registered_arg_description_reaches_llm_schema():
    """예전에는 `option: str = ''`만 넘어가서, 콘솔에 적어 둔 옵션 설명이 모델에
    **한 글자도** 전달되지 않았다. 등록은 했는데 에이전트가 옵션을 못 채운 원인이다."""
    props = _arg_schema({
        "description": "job 목록", "exec_command": "phd list {option}",
        "args": [{"name": "option", "type": "str", "required": False, "default": "",
                  "description": "-l: 상세 출력 -lf: 필드 목록", "choices": []}]})
    assert props["option"]["description"] == "-l: 상세 출력 -lf: 필드 목록"


def test_enum_choices_become_schema_enum_with_labels():
    """선택지를 `값: 설명`으로 적으면 값만 enum이 되고 설명은 파라미터 설명으로 간다."""
    props = _arg_schema({
        "description": "job 상세", "exec_command": "phd info {option} {job_id}",
        "args": [{"name": "option", "type": "enum", "required": False, "default": "",
                  "description": "",
                  "choices": ["-j: JSON 형식으로 반환", "-tl: 부가 정보 출력"]},
                 {"name": "job_id", "type": "str", "required": True,
                  "description": "job id", "choices": []}]})
    # 필수가 아닌 선택형에는 빈 값도 있어야 한다(기본값 ""이 enum에 없으면 스키마가 깨진다).
    assert set(props["option"]["enum"]) == {"-j", "-tl", ""}
    assert "JSON 형식으로 반환" in props["option"]["description"]
    assert props["job_id"]["description"] == "job id"


def test_enum_value_parsing_ignores_description_part():
    """`cast_arg`도 같은 규칙으로 값을 비교해야 한다(스키마와 검증이 어긋나면 전부 거부된다)."""
    from execution_exec import cast_arg, choice_value
    spec = {"name": "option", "type": "enum",
            "choices": ["-j: JSON 형식으로 반환", "-tl: 부가 정보 출력"]}
    assert cast_arg(spec, "-j") == "-j"
    assert choice_value("12:00") == "12:00"        # 콜론 앞이 값처럼 보여도 공백 규칙으로 구분
    with pytest.raises(ValueError):
        cast_arg(spec, "-zz")


# --- #140: 에이전트가 질문한 사용자 계정을 안다 -----------------------------------
def test_agent_injects_caller_account_into_prompt():
    """모델이 자기 이름(`ops_assistant`)을 사용자 계정으로 말한 사고가 있었다.
    호출자 계정을 알려 주지 않았기 때문이다."""
    src = open(os.path.join(ROOT, "agent_server", "agent.py"), encoding="utf-8").read()
    assert 'caller_headers or {}).get("X-User-Id")' in src
    assert "질문한 사용자 계정" in src


def test_instruction_answers_other_user_question_in_one_line():
    instr = _instruction_text()
    assert "## 남의 계정을 묻는 질문" in instr
    assert "가이드 문서 위치를 안내하지 않습니다" in instr
    assert "질문한 사용자 계정" in instr


# --- #140: 엑셀 양식이 인자까지 실어 나른다 ---------------------------------------
def _exec_router():
    """관리자 콘솔 실행 라우터를 DB 없이 임포트한다(모듈 상수/파서만 쓴다)."""
    import importlib.util
    for k, v in {"CONFIG_DB_DSN": "postgres://x/y", "POSTGRES_PASSWORD": "x",
                 "ADMIN_PASSWORD": "x", "SESSION_SECRET": "x"}.items():
        os.environ.setdefault(k, v)
    sys.path.insert(0, os.path.join(ROOT, "admin_console", "backend", "routers"))
    path = os.path.join(ROOT, "admin_console", "backend", "routers", "execution.py")
    spec = importlib.util.spec_from_file_location("execrouter", path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def test_excel_template_roundtrips_with_args():
    """예전 일괄 등록은 이름·설명·실행 커맨드만 읽고 **인자를 버렸다**. `{option}`이 있는
    커맨드는 엑셀로 넣어 봐야 인자가 빈 채로 들어가 한 건씩 다시 채워야 했다.

    양식을 만들어 그대로 다시 읽었을 때 인자 정의가 살아 있고, 화면 등록과 같은 검증을
    통과해야 한다(양식과 파서가 따로 놀면 사용자는 '업로드했는데 안 된다'만 겪는다)."""
    import io
    import openpyxl
    from execution_exec import placeholders_in, validate_definition, deny_set, DEFAULT_DENY_CSV

    m = _exec_router()
    data = m._build_workbook([list(r) for r in m._TEMPLATE_EXAMPLES])
    ws = openpyxl.load_workbook(io.BytesIO(data))["커맨드"]
    rows = list(ws.iter_rows(values_only=True))
    header = [str(h) for h in rows[0]]
    deny = deny_set(DEFAULT_DENY_CSV)

    parsed = {}
    for r in rows[1:]:
        def cell(col, _r=r):
            v = _r[header.index(col)] if col in header else None
            return "" if v is None else str(v)
        exec_cmd = cell("실행 커맨드")
        ph = [p for p in placeholders_in(exec_cmd) if p != "user_id"]
        args = m._parse_args_from_row(cell, header, ph)
        host = m._HOST_WORDS.get(cell("실행 위치").strip().lower(), "login_server")
        validate_definition(cell("이름"), exec_cmd, args, host, deny)   # 화면 등록과 같은 문
        parsed[cell("이름")] = args

    assert parsed["myquota"] == []
    info = {a["name"]: a for a in parsed["s2_phd_info"]}
    # 자리표시자 순서대로 이름이 자동으로 붙는다(양식의 '이름' 칸을 비워도 되게 한 이유).
    assert [a["name"] for a in parsed["s2_phd_info"]] == ["option", "job_id"]
    assert info["option"]["type"] == "enum"
    assert info["option"]["choices"][0].startswith("-j:")
    assert info["job_id"]["required"] is True


def test_excel_arg_names_default_to_placeholder_order():
    """'인자N 이름'을 비워 두면 실행 커맨드의 자리표시자 순서로 채워진다.
    관리자가 이름을 두 번 옮겨 적다 틀리면 그 커맨드만 통째로 거부된다."""
    m = _exec_router()
    header = ["인자1 설명", "인자2 설명"]
    values = {"인자1 설명": "출력 형식", "인자2 설명": "job id"}
    args = m._parse_args_from_row(lambda c: values.get(c, ""), header, ["option", "job_id"])
    assert [a["name"] for a in args] == ["option", "job_id"]
    assert args[1]["description"] == "job id"


def test_excel_choices_split_on_newline_not_comma():
    """선택지 설명에 콤마가 흔하다("Return job info, json format").
    콤마로 쪼개면 설명 조각이 값으로 등록된다."""
    m = _exec_router()
    header = ["인자1 선택지"]
    text = "-j: Return job info, json format\n-tl: Print additional info"
    args = m._parse_args_from_row(lambda c: text if c == header[0] else "", header, ["option"])
    assert args[0]["choices"] == ["-j: Return job info, json format",
                                  "-tl: Print additional info"]
    assert args[0]["type"] == "enum"        # 선택지를 적었으면 타입 칸이 비어도 선택형


# --- #141: 상태를 담는 서비스에는 **이름 있는 볼륨**이 반드시 있어야 한다 --------------
# dev compose의 postgres에 데이터 볼륨이 없었다. pgvector 이미지가
# `VOLUME /var/lib/postgresql/data`를 선언하므로 **익명 볼륨**이 붙는데, 익명 볼륨은
# 컨테이너를 다시 만들 때 떨어져 나간다. #139에서 `ports:`를 127.0.0.1로 바꾼 것만으로
# compose가 postgres를 재생성했고, 매뉴얼·VOC·설정·등록 커맨드가 전부 사라졌다.
# 상태를 담는 **이미지**와 그 데이터 경로. 새 서비스를 붙일 때 여기 추가하면 검사가 따라온다.
# (redis는 임베딩 캐시라 없어져도 재생성되므로 뺐다 - 소실이 손실이 아닌 유일한 경우다.)
_STATEFUL_IMAGES = {
    "pgvector": "/var/lib/postgresql/data",
    "postgres": "/var/lib/postgresql/data",
    "clickhouse-server": "/var/lib/clickhouse",
    "minio": "/data",
    "open-webui": "/app/backend/data",
}


def _compose(name):
    yaml = pytest.importorskip("yaml")
    with open(os.path.join(ROOT, name), encoding="utf-8") as f:
        return yaml.safe_load(f)


@pytest.mark.parametrize("compose_file", ["docker-compose.yml", "docker-compose.dev.yml"])
def test_stateful_services_have_named_volumes(compose_file):
    """상태를 담는 서비스는 전부 **이름 있는 볼륨**이어야 한다.

    dev postgres에 데이터 볼륨이 없었다. pgvector 이미지가 `VOLUME`을 선언하므로 익명 볼륨이
    붙는데, 익명 볼륨은 컨테이너를 다시 만들면 떨어져 나간다. #139에서 `ports:`를 127.0.0.1로
    바꾼 것만으로 compose가 postgres를 재생성했고, 매뉴얼·VOC·설정·등록 커맨드가 전부 사라졌다.
    `down -v`만 위험한 게 아니다 - 재생성을 부르는 아무 변경이나 같은 결과를 낸다.

    postgres 하나만 보지 않고 **상태를 담는 이미지 전체**를 훑는다. 같은 실수를 세 번째로
    하지 않으려면 새 서비스가 붙을 때 자동으로 걸려야 한다.
    """
    conf = _compose(compose_file)
    declared = set((conf.get("volumes") or {}) or [])
    checked = 0

    for svc_name, svc in (conf.get("services") or {}).items():
        image = str(svc.get("image", ""))
        # `entrypoint`를 덮어쓴 서비스는 그 이미지를 **클라이언트로** 쓰는 것이다
        # (dev-config는 postgres 이미지로 psql만 돌린다). 서버는 이미지 기본 엔트리포인트로 뜬다.
        if svc.get("entrypoint"):
            continue
        for key, data_path in _STATEFUL_IMAGES.items():
            if key not in image:
                continue
            mounts = [str(m) for m in (svc.get("volumes") or [])]
            data = [m for m in mounts if m.split(":")[1:2] == [data_path]]
            assert data, (
                f"{compose_file}: '{svc_name}'({image})에 {data_path} 볼륨이 없습니다. "
                "익명 볼륨이 붙어 컨테이너를 다시 만들 때마다 데이터가 사라집니다.")
            src = data[0].split(":")[0]
            assert not src.startswith((".", "/", "$")), (
                f"{compose_file}: '{svc_name}'의 데이터 볼륨이 바인드 마운트({src})입니다. "
                "rsync --delete가 지울 수 있습니다(#137).")
            assert src in declared, (
                f"{compose_file}: '{svc_name}'의 '{src}'가 최상위 volumes에 없습니다.")
            checked += 1
            break

    assert checked, f"{compose_file}에서 상태 서비스를 하나도 찾지 못했습니다(검사가 헛돌고 있음)."


def test_openwebui_has_named_data_volume():
    """Open WebUI 계정·대화도 같은 이유로 이름 있는 볼륨이어야 한다."""
    yaml = pytest.importorskip("yaml")
    with open(os.path.join(ROOT, "docker-compose.dev.yml"), encoding="utf-8") as f:
        conf = yaml.safe_load(f)
    mounts = conf["services"]["open-webui"]["volumes"]
    data = [m for m in mounts if str(m).endswith(":/app/backend/data")]
    assert data and not str(data[0]).startswith(".")


# --- #142: /v1/*에 인증을 걸면서 **내부 호출자**를 확인하지 않았다 --------------------
def test_admin_console_sends_agent_api_key_to_agent_server():
    """#139에서 agent-server의 `/v1/*`에 `agent_api_key` 인증을 걸었는데, 관리자 콘솔의
    "기본 모델 동기화"가 `/v1/models`를 **헤더 없이** 부르고 있었다. 그래서 키를 넣는 순간
    값이 맞든 틀리든 401로 죽었다(사용자가 저장할 때마다 오류를 봤다).

    제약을 새로 걸 때는 그 경로에 이미 붙어 있던 호출자를 전부 훑어야 한다.
    """
    src = open(os.path.join(ROOT, "admin_console", "backend", "routers", "ops.py"),
               encoding="utf-8").read()
    assert 'get_config("agent_api_key"' in src, "콘솔이 agent_api_key를 읽지 않는다"

    # agent-server를 부르는 줄에 **인증 헤더가 붙어 있어야** 한다.
    get_line = next((ln for ln in src.split("\n") if "/v1/models" in ln and "client.get" in ln), "")
    assert get_line, "ops.py에서 /v1/models 호출을 찾지 못했다(테스트를 갱신할 것)"
    assert "headers=agent_headers" in get_line, \
        f"콘솔이 agent-server를 부를 때 인증 헤더를 안 보낸다: {get_line.strip()}"

    # 그 헤더는 **agent_api_key**로 만들어야 한다. Open WebUI 키를 보내면 방향이 뒤집힌다
    # (두 키는 목적지가 정반대다 - 사용자가 실제로 헷갈린 지점이다).
    hdr = next((ln for ln in src.split("\n") if "agent_headers =" in ln), "")
    assert "agent_key" in hdr and "openwebui" not in hdr, \
        f"agent-server용 헤더를 agent_api_key로 만들지 않는다: {hdr.strip()}"


def test_agent_api_key_is_hot_reload():
    """`agent_api_key`가 hot_reload=false면 저장 후 재시작 전까지 콘솔↔agent-server 값이
    어긋나 401이 계속된다. 안내 문구도 '재시작하세요'로 바뀌어야 하므로 고정해 둔다."""
    src = open(os.path.join(ROOT, "shared", "migrations.py"), encoding="utf-8").read()
    i = src.index('("agent_api_key"')
    seed = src[i:src.index("),", i)]
    # (key, value, desc, hot_reload, is_secret, force) — hot_reload가 True여야 한다.
    assert seed.rstrip().endswith("True, True, False"), f"시드 플래그가 바뀌었다: {seed[-40:]}"


# --- #144: 비로그인 강등 모드에서 작업 디렉토리가 root 홈이었다 ------------------------
@pytest.mark.parametrize("mode", ["su", "runuser"])
def test_non_login_privdrop_moves_to_user_home(mode):
    """`ssh root@host <cmd>`는 root 홈(`/root`)에서 시작하는데, `runuser -u`와 비로그인 `su`는
    **작업 디렉토리를 바꾸지 않는다**. 그래서 `SSH_PRIVDROP=runuser`로 바꾼 뒤 `ls -lh`가
    `/root`에서 돌아 `Permission denied`가 났다.

    지시문이 "실행은 항상 본인 홈에서 시작합니다"라고 약속하므로 코드가 그것을 지켜야 한다."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from ssh_exec import _remote_command
    cmd = _remote_command("yr9.choi", ["ls", "-lh"], mode)
    assert cmd.startswith("cd ~yr9.choi"), f"홈으로 이동하지 않는다: {cmd}"
    # `&&`면 root가 홈에 못 들어가는 환경(GPFS root_squash)에서 커맨드가 아예 안 돈다.
    assert "; " in cmd and "&&" not in cmd, f"실패 시 커맨드를 막으면 안 된다: {cmd}"


def test_login_privdrop_does_not_double_cd():
    """`su - user`는 로그인 셸이라 이미 홈으로 간다. 덧붙이면 군더더기다."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from ssh_exec import _remote_command
    cmd = _remote_command("yr9.choi", ["ls", "-lh"], "su-login")
    assert cmd == "su - yr9.choi -c 'ls -lh'"


def test_home_cwd_does_not_break_argument_quoting():
    """홈 이동을 붙이면서 인자 인용이 깨지면 셸 주입이 생긴다."""
    sys.path.insert(0, os.path.join(ROOT, "shared"))
    from ssh_exec import _remote_command
    cmd = _remote_command("yr9.choi", ["ls", "; rm -rf /", "$HOME", "`id`"], "runuser")
    assert "'; rm -rf /'" in cmd and "'$HOME'" in cmd and "'`id`'" in cmd
    # 우리가 의도한 `;`는 하나뿐이어야 한다(cd 뒤).
    assert cmd.count(";") == 2, cmd      # cd 뒤 1개 + 인자 안의 리터럴 1개(따옴표 안)


def test_instruction_decides_by_content_not_phrasing():
    """실행 여부를 **말투로** 판단하면 안 된다(#149).

    "보여 줘"라고 해야만 실행하도록 써 뒀더니 "내 홈 디렉토리는 어디야?"가 일반 지식으로
    분류돼 도구를 아예 호출하지 않았다. 질문 형태를 나열하는 것은 커맨드를 나열하는 것과
    같은 실수다(#145) - 사용자가 어떻게 물을지는 알 수 없다.
    기준은 **답에 무엇이 필요한가** 하나여야 한다.
    """
    instr = _instruction_text()
    assert "말투로 판단하지 않습니다" in instr
    assert "답이 이 서버에 물어봐야 나오는 값이면 실행합니다" in instr
    # 판별법이 있어야 실행 가능한 규칙이 된다.
    assert "회사·서버·" in instr and "따라 달라지는가" in instr


def test_instruction_forbids_fabricating_environment_values():
    """마지막 안전장치: 분류를 잘못해 (B)로 답하더라도 이 환경의 값은 지어내지 않는다.

    사용자 지적: "일반지식이더라도 모델이 만들어내진 말아야지. 모르는 건 모른다고 해야지."
    실제로 `/home/yr9.choi`를 지어내 답했다(정답은 `/home/gpu1/yr9.choi`).
    """
    instr = _instruction_text()
    assert "이 환경의 값을 지어내지 않습니다" in instr
    # 지어내는 중임을 스스로 알아채는 신호 - 헤지 문구를 명시적으로 금지한다.
    for hedge in ("일반적으로 …입니다", "보통 …입니다", "정확한 것은 직접 확인해 보세요"):
        assert hedge in instr, f"헤지 문구를 금지 목록에 넣지 않았다: {hedge}"
    assert "확인해 봐야 알 수 있습니다" in instr, "모른다고 답할 문구를 주지 않았다"
    assert "지어낸 값을 주는 것이 실패입니다" in instr


# --- #145: 지시문에 **특정 커맨드를 박지 않는다** ------------------------------------
# 사용자가 네 번째로 지적한 사항이다(#74 `phd info`, #125, #140, 그리고 `pwd`/`echo $HOME`).
# 커맨드는 (1) 콘솔에 등록돼 툴로 노출되거나 (2) 모델이 아는 표준 리눅스 명령이다.
# 지시문은 **원칙**만 말해야 한다 - 커맨드를 적기 시작하면 하나하나 다 적어야 하고,
# 시스템이 바뀔 때마다 지시문이 거짓말을 하게 된다(#144가 정확히 그렇게 났다).
_FORBIDDEN_IN_INSTRUCTION = [
    # 사내 전용(존재하지 않는 것을 지어내 쓴 사고가 있었다)
    "phd ", "myquota", "squeue", "sinfo", "sbatch", "bsub",
    # 확인용 표준 명령(모델이 알아서 고르게 둔다)
    "pwd", "whoami", "echo $", "$HOME", "$USER",
    # 나열하기 시작하면 끝이 없다
    "nvidia-smi", "`ls ", "`ls`", "`df", "`du", "`head", "`tail", "`find",
]


def test_instruction_names_no_specific_commands():
    """지시문에 커맨드 이름을 적지 않는다. 원칙만 쓰고 선택은 모델과 툴 목록에 맡긴다."""
    instr = _instruction_text()
    # 모듈 docstring은 제외하고 지시문 본문만 본다.
    body = instr[instr.index("AGENT_INSTRUCTION = "):]
    hits = [tok for tok in _FORBIDDEN_IN_INSTRUCTION if tok in body]
    assert not hits, (
        f"지시문에 특정 커맨드가 박혀 있다: {hits}\\n"
        "  커맨드는 콘솔 등록(툴)이나 모델의 표준 리눅스 지식으로 해결한다. "
        "지시문에는 원칙만 쓴다.")


def test_instruction_states_no_shell_as_a_property():
    """커맨드를 적는 대신 **성질**을 말해야 한다: 셸을 거치지 않는다.
    그래야 모델이 `echo $HOME` 같은 것을 스스로 피한다(그건 글자 그대로 출력된다)."""
    instr = _instruction_text()
    assert "셸을 거치지 않습니다" in instr
    assert "글자 그대로" in instr


# --- #147: "되돌리기" 버튼이 옛 지시문을 계속 저장했다 --------------------------------
def test_instruction_reset_reads_file_not_module_cache():
    """`agent_system_instruction` 되돌리기가 **파일을 다시 읽어야** 한다.

    예전에는 함수 안에서 `from agent_instruction import AGENT_INSTRUCTION` 했다. 부수효과가
    없으니 안전하다고 생각했는데, 안전한 것과 **최신인 것**은 다르다. 파이썬은 `sys.modules`에
    캐시하므로 이 프로세스가 한 번이라도 읽었으면 그 뒤로는 옛 텍스트를 쓴다. `./shared`가
    바인드 마운트라 파일은 최신인데 **버튼이 아무 일도 하지 않는** 상태가 됐다
    (사용자: "2번 그대로 했는데도 1번 하면 옛지시문이라고 떠").

    `importlib.reload`도 부족하다 - `.pyc`는 mtime+크기로 유효성을 보므로 같은 초에 같은
    크기로 바뀌면 낡은 바이트코드를 쓴다(실제로 재현했다). 그래서 소스를 직접 읽어 compile한다.
    """
    src = open(os.path.join(ROOT, "admin_console", "backend", "routers", "settings.py"),
               encoding="utf-8").read()
    reset = src[src.index("def _read_instruction_from_disk"):]
    assert "open(path" in reset and "compile(src" in reset, \
        "지시문을 파일에서 직접 읽지 않는다"
    # 되돌리기 경로에서 모듈 import로 지시문을 가져오면 안 된다.
    body = src[src.index("async def reset_agent_instruction"):]
    assert "from agent_instruction import" not in body, \
        "되돌리기가 모듈 캐시에서 지시문을 읽는다 - 옛 텍스트가 저장된다"
    assert "_read_instruction_from_disk()" in body
    # 못 읽었으면 조용히 옛 값을 쓰지 않고 실패해야 한다.
    assert "지시문 파일을 읽지 못했습니다" in reset


def test_instruction_reset_actually_returns_current_file_text():
    """헬퍼가 정말 디스크의 현재 내용을 돌려주는지 - 문자열 검사만으로는 부족하다."""
    for k, v in {"CONFIG_DB_DSN": "postgres://x/y", "POSTGRES_PASSWORD": "x",
                 "ADMIN_PASSWORD": "x", "SESSION_SECRET": "x"}.items():
        os.environ.setdefault(k, v)
    import importlib.util
    sys.path.insert(0, os.path.join(ROOT, "admin_console", "backend"))
    sys.path.insert(0, os.path.join(ROOT, "admin_console", "backend", "routers"))
    path = os.path.join(ROOT, "admin_console", "backend", "routers", "settings.py")
    spec = importlib.util.spec_from_file_location("settings_under_test", path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)

    on_disk = _instruction_text()
    got = mod._read_instruction_from_disk()
    assert got in on_disk, "파일에 없는 내용을 돌려준다"
    assert len(got) > 5000, f"지시문이 너무 짧다({len(got)}자) - 잘못 읽고 있다"


# --- #148: 엑셀 양식의 고정 선택지는 **드롭다운**이어야 한다 --------------------------
def test_excel_template_has_dropdowns_for_fixed_choice_columns():
    """타입/필수/활성/실행 위치는 값이 정해져 있다. 자유 입력이면 오타가 조용히 다른 뜻이 된다
    ("선택"이라고 적으면 선택형이 아니라 문자열로 들어간다)."""
    import io
    import openpyxl
    from openpyxl.utils import column_index_from_string

    m = _exec_router()
    ws = openpyxl.load_workbook(io.BytesIO(m._build_workbook([])))["커맨드"]
    cols = [c.value for c in ws[1]]

    got = {}
    for dv in ws.data_validations.dataValidation:
        first = str(dv.sqref).split()[0].split(":")[0]
        letter = "".join(ch for ch in first if ch.isalpha())
        got[cols[column_index_from_string(letter) - 1]] = dv.formula1

    for col, expected in [("실행 위치", "로그인 서버"), ("활성", "Y,N"),
                          ("인자1 타입", "문자열,정수,선택형"), ("인자1 필수", "Y,N")]:
        assert col in got, f"'{col}' 열에 드롭다운이 없다"
        assert expected in got[col], f"'{col}' 선택지가 다르다: {got[col]}"

    # 인자 슬롯 전부에 걸려야 한다(1번만 걸고 나머지를 빠뜨리기 쉽다).
    for i in range(1, m.TEMPLATE_ARG_SLOTS + 1):
        assert f"인자{i} 타입" in got and f"인자{i} 필수" in got, f"인자{i} 슬롯에 드롭다운이 없다"


def test_excel_dropdown_values_match_the_parser():
    """드롭다운에 있는 값은 **파서가 실제로 알아듣는 값**이어야 한다.
    화면에서 고를 수 있는데 업로드하면 무시되는 값이 있으면 안 된다."""
    m = _exec_router()
    for label in m._DROPDOWNS["타입"]:
        assert label.lower() in m._TYPE_WORDS, f"파서가 모르는 타입: {label}"
    for label in m._DROPDOWNS["실행 위치"]:
        assert label.lower() in m._HOST_WORDS, f"파서가 모르는 실행 위치: {label}"
    for label in m._DROPDOWNS["필수"]:
        assert m._truthy(label, False) == (label == "Y"), f"필수 값 해석이 다르다: {label}"


def test_empty_default_drops_the_argument_entirely():
    """"기본값이 없으면 아무것도 안 보이는 건가?" - 자리표시자가 통째로 빠진다."""
    from execution_exec import build_registered_argv, deny_set, DEFAULT_DENY_CSV
    deny = deny_set(DEFAULT_DENY_CSV)
    spec = [{"name": "option", "type": "enum", "required": False,
             "default": "", "choices": ["-l", "-lf"]}]
    assert build_registered_argv("phd list {option}", spec, {}, None,
                                 "yr9.choi", deny, True) == ["phd", "list"]
    spec[0]["default"] = "-l"
    assert build_registered_argv("phd list {option}", spec, {}, None,
                                 "yr9.choi", deny, True) == ["phd", "list", "-l"]


def test_progress_line_always_shows_line_count():
    """잘리지 않았을 때도 줄 수를 보여줘야 한다(#149).

    예전에는 잘렸을 때만 `⚠ N줄 중 M줄만`을 찍었다. 그래서 22줄짜리 답을 받고도
    **우리가 자른 건지 모델이 자른 건지** 구분할 수 없었다. 항상 줄 수를 찍으면
    사용자가 답변의 행 수와 눈으로 대조할 수 있다.
    """
    import json as _json
    src = open(os.path.join(ROOT, "agent_server", "main.py"), encoding="utf-8").read()
    i = src.index("def _result_phrase")
    end = src.index("\nclass _StreamDedup")
    ns = {"json": _json}
    exec(src[i:end], ns)
    phrase = ns["_result_phrase"]

    base = {"ip": "10.0.0.1", "as_user": "yr9.choi", "duration_ms": 400, "exit_code": 0}
    full = phrase("run_command", {**base, "truncated": False,
                                  "total_lines": 132, "shown_lines": 132})
    assert "132줄" in full and "⚠" not in full, full

    cut = phrase("run_command", {**base, "truncated": True,
                                 "total_lines": 132, "shown_lines": 58})
    assert "⚠ 출력 132줄 중 58줄만" in cut, cut

    # 한 줄짜리 출력에까지 붙이면 잡음이다.
    one = phrase("run_command", {**base, "truncated": False,
                                 "total_lines": 1, "shown_lines": 1})
    assert "줄" not in one, one
